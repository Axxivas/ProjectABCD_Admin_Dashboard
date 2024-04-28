import sys, io, requests, threading, webbrowser, urllib, os, platform, openpyxl, string, textwrap, re, time, textstat, datetime, os,pyttsx3,random,tempfile,textstat,traceback, openai
import pandas as pd
import tkinter as tk
from collections import defaultdict
from concurrent.futures import ThreadPoolExecutor, as_completed
from gtts import gTTS
from localspelling import convert_spelling
from localspelling.spelling_converter import get_dictionary
from pandas import DataFrame, Series
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.shapes.connector import Connector
from pptx.util import Inches, Pt
from textblob import TextBlob
from tkinter import filedialog, messagebox, ttk
from xml.dom.expatbuilder import FragmentBuilderNS
import googletrans

openai.api_key = 'private'

ROOT_WIDTH = 1000 # app window width
ROOT_HEIGHT = 600 # app window height

root = tk.Tk()
root.title("Project ABCD Admin Panel")
sw_placement = int(root.winfo_screenwidth()/2 - ROOT_WIDTH/2) # to place at half width of screen
sh_placement = int(root.winfo_screenheight()/2 - ROOT_HEIGHT/2) # to place at half height of screen
root.geometry(f"{ROOT_WIDTH}x{ROOT_HEIGHT}+{sw_placement}+{sh_placement}")
root.minsize(ROOT_WIDTH, ROOT_HEIGHT)

# Set sys.stdout to use utf-8 encoding
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

MAIN_FONT = ("helvetica", 12)
LABEL_FONT = ("helvetica bold", 14)
LANGUAGES = [
    "Telugu",
    "Hindi",
    "Spanish"
]

# creates preferences dictionary from preferences.txt
# creates default preferences.txt file if one does not exist
try:
    with open("preferences.txt", "r", encoding="utf8") as file:
        lines = file.readlines()
        preferences = {}

        for line in lines:
            key, value = line.split('=')
            preferences[key.strip()] = value.strip().replace('“', '').replace('”', '').replace('"', '').replace("'", '')
except FileNotFoundError:
    preferences = {
        "TEXT_SIZE" : "14",
        "TEXT_FONT" : "Times New Roman",
        "TITLE_SIZE" : "32",
        "TITLE_FONT" : "Arial",
        "SUBTITLE_SIZE" : "24",
        "SUBTITLE_FONT" : "Arial",
        "PIC_WIDTH" : "720",
        "PIC_HEIGHT" : "1040"
    }
    tk.messagebox.showwarning(title='Warning', message='No preferences.txt file exists in directory. Default preferences.txt will be created and used.')
    print('No preferences.txt file exists in directory. Default preferences.txt will be created and used.')
    with open('preferences.txt', 'w') as f:
        f.write('TEXT_SIZE = 14\nTEXT_FONT = Times New Roman\nTITLE_SIZE = 32\nTITLE_FONT = Arial\nSUBTITLE_SIZE = 24\nSUBTITLE_FONT = Arial\nPIC_WIDTH = 720\nPIC_HEIGHT = 1040')

'''
Generates an xcel file with the pairs found
'''
def generatePairs():
    from origPythonCode import apiRunner, generate_table
    file_path = 'APIData.xlsx'  # Change to the actual file path
    api_dress_data = sorted(apiRunner(), key=lambda x: x['id'])  # Assuming apiRunner() returns dress data
    pairs = []

    try:
        # Read dress data from Excel file
        sheet_dress_data = pd.read_excel(file_path)      
        sheet_dress_data.dropna(subset=['id'], inplace=True)  # Drop rows with missing IDs
        sheet_dress_data['description'].fillna('', inplace=True)  # Remove NA/NAN from description column
        sheet_dress_data['did_you_know'].fillna('', inplace=True)  # Remove NA/NAN from did_you_know column

	# Iterate through API dress data
        for data_that_will_be_searched in api_dress_data:
            # Get the name of the current API data ID
            name = data_that_will_be_searched['name']
            # Split the name into tokens and remove any prefixes
            tokens = [token + " " for token in name.split() if token not in ["Dr.", "Mr.", "Mrs.", "Ms."]]

            # Loop through each token
            for token in tokens:
                # Loop through provided IDs instead of the entire sheet_dress_data
                for data_that_contains_token in api_dress_data:
                    # Get the description and did you know text of the provided ID
                    description = data_that_contains_token['description']
                    did_you_know = data_that_contains_token['did_you_know']

                    # Check if the token is present in either of them, make sure we aren't looking on the same IDs
                    if data_that_will_be_searched['id'] != data_that_contains_token['id']:
                        if token in description or token in did_you_know:
                            if (data_that_contains_token['id'], data_that_will_be_searched['id']) not in [(pair[0], pair[2]) for pair in pairs]:
                                # Add the pair of IDs and names to the list
                                pairs.append([data_that_contains_token['id'], data_that_contains_token['name'], data_that_will_be_searched['id'], name])
                            # Break the inner loop as we found a pair for this token
                            break

        # Generate table and save to Excel
        column_headers = ['ID1', 'Name 1', 'ID2', 'Name 2']
        generate_table(pairs, 'generate_pairs', column_headers, 50, 200, 'center', 1)
        df_pairs = pd.DataFrame(pairs, columns=column_headers)
        df_pairs.to_excel("pairs_generated.xlsx", index=False)
        print("Excel file 'pairs_generated.xlsx' created.")

    except FileNotFoundError:
        print(f"File '{file_path}' not found.")
    except Exception as e:
        print(f'Error: {e}')

def generate_us_uk_spellings():
    from origPythonCode import apiRunner,generate_table
    file_path = 'APIData.xlsx'  # Change to the actual file path
    api_dress_data = sorted(apiRunner(), key=lambda x: x['id'])  # Assuming apiRunner() returns dress data
    words = []

    try:
        # Read dress data from Excel file
        sheet_dress_data = pd.read_excel(file_path)      
        sheet_dress_data.dropna(subset=['id'], inplace=True)  # Drop rows with missing IDs
        sheet_dress_data['description'].fillna('', inplace=True)  # Remove NA/NAN from description column
        sheet_dress_data['did_you_know'].fillna('', inplace=True)  # Remove NA/NAN from did_you_know column

	# Iterate through API dress data
        for api_data in api_dress_data:
            ID = api_data['id']
            name = api_data['name']
            description = api_data['description']
            did_you_know = api_data['did_you_know']

            # concatenation of description and did_you_know
            merged_words = f'{str(description)} {str(did_you_know)}'
            blob = TextBlob(merged_words)

            uk_words = []
            us_words = []
            uk_version_of_the_us_words = []
            same_words_for_both = []

            # Iterate over each word in the text
            for word in blob.words:
                # to check if it is us or uk
                if word in get_dictionary("us").values():
                    us_words.append(word)
                    uk_version_of_the_us_words.append(convert_spelling(word, 'gb'))

                elif word in get_dictionary("gb").values():
                    uk_words.append(word)
                else:
                    same_words_for_both.append(word)     

            words.append([ID, name, us_words, uk_version_of_the_us_words])

        # Generate table and save to Excel
        column_headers = ['ID', 'Name', 'US spellings', 'UK spellings']
        generate_table(words, 'US_UK Spellings', column_headers, 50, 300, 'center', 1)
        df_pairs = pd.DataFrame(words, columns=column_headers)
        df_pairs.to_excel("US_UK_Spelling.xlsx", index=False)
        print("Excel file 'US_UK_Spelling.xlsx' created.")

    except FileNotFoundError:
        print(f"File '{file_path}' not found.")
    except Exception as e:
        print(f'Error: {e}')

# ids_of_us_uk spelling function
def generate_IDs_Of_us_uk_spellings():
    from origPythonCode import apiRunner, generate_table

    from collections import defaultdict

    file_path = 'APIData.xlsx'  # Change to the actual file path
    api_dress_data = sorted(apiRunner(), key=lambda x: x['id'])  # Assuming apiRunner() returns dress data
    # Initialize a defaultdict to store dress IDs for each US spelling
    us_spellings_dict = defaultdict(lambda: {'ids': []})

    try:
        # Read dress data from Excel file
        sheet_dress_data = pd.read_excel(file_path)      
        sheet_dress_data.dropna(subset=['id'], inplace=True)  # Drop rows with missing IDs
        sheet_dress_data['description'].fillna('', inplace=True)  # Remove NA/NAN from description column
        sheet_dress_data['did_you_know'].fillna('', inplace=True)  # Remove NA/NAN from did_you_know column

        unique_ids_words_pairs = []  # Initialize an empty list to store the lists

        for api_data in api_dress_data:
            ID = api_data['id']
            description = api_data['description']
            did_you_know = api_data['did_you_know']

            # Concatenation of description and did_you_know
            merged_words = f'{str(description)} {str(did_you_know)}'
            blob = TextBlob(merged_words)

            uk_words = []
            us_words = []
            uk_version_of_the_us_words = []
            same_words_for_both = []

            # Iterate over each word in the text
            for word in blob.words:
                # to check if it is us or uk
                if word in get_dictionary("us").values():
                    us_words.append(word)
                    uk_version_of_the_us_words.append(convert_spelling(word, 'gb'))

                    # Append dress ID to the corresponding US spelling
                    us_spellings_dict[word]['ids'].append(ID)

                elif word in get_dictionary("gb").values():
                    uk_words.append(word)
                else:
                    same_words_for_both.append(word)

        # Create the list of unique_ids_words_pairs
        for us_word, data in us_spellings_dict.items():
            ids = data['ids']
            unique_ids = list(set(ids))  # Convert to set to remove duplicates, then back to list
            uk_word = convert_spelling(us_word, 'gb')
            
            # Create a new list with the US word, UK spelling, and unique IDs
            pair = [us_word, uk_word, unique_ids]
            
            # Append the list to the list
            unique_ids_words_pairs.append(pair)

        # Generate table and save to Excel
        column_headers = ['US', 'UK', 'IDS']
        generate_table(unique_ids_words_pairs, 'IDS_OF_US_UK Spellings', column_headers, 50, 300, 'center', 1)
        df_pairs = pd.DataFrame(unique_ids_words_pairs, columns=column_headers)
        df_pairs.to_excel("IDS_OF_US_UK Spellings.xlsx", index=False)
        print("Excel file 'IDS_OF_US_UK Spellings.xlsx' created.")

    except FileNotFoundError:
        print(f"File '{file_path}' not found.")
    except Exception as e:
        print(f'Error: {e}')

'''
Function to fetch the english text from dress data
'''
def fetch_english_text(dress_data):
    english_texts = {}  
    for dress in dress_data:
        dress_id = dress.get('id')
        english_description = {
            'name': dress.get('name', ''),
            'description': dress.get('description', ''),
            'did_you_know': dress.get('did_you_know', '')
        }
        english_texts[dress_id] = english_description
    return english_texts
'''
translation function that stores the translated telugu texts into the same "keys" as the english for easier formatting
'''
def translate_text_to_telugu(english_texts):
    translator = googletrans.Translator()
    telugu_texts = {}

    for id, texts in english_texts.items():
        telugu_texts[id] = {}  # Initialize a dictionary for this ID
        for key, text in texts.items():
            try:
                translated = translator.translate(text, dest='te')  # 'te' for Telugu
                telugu_texts[id][key] = translated.text
            except Exception as e:
                print(f"Error translating {key} for ID {id}: {e}")
                telugu_texts[id][key] = text  # Use the original text if translation fails

    return telugu_texts




def read_translated_ids(filename):
    translated_ids = set()
    if os.path.exists(filename):
        with open(filename, 'r') as file:
            for line in file:
                translated_ids.add(line.strip())
    return translated_ids

def write_translated_ids(translated_ids, filename):
    with open(filename, 'w') as file:
        for id in translated_ids:
            file.write(f"{id}\n")

'''
Performs change of third person text to first person
'''
def translate_text_to_first_person(english_texts):
    i = 1
    tries = 1
    sleep_timer = 0
    first_person_texts = {} # Dictionary to hold the translated text
    translated_ids_file = 'translated_ids.txt'
    translated_ids = read_translated_ids(translated_ids_file)
    print(translated_ids, flush=True)


    with open("output.txt", "a") as file: # Opens a text file to write the translated text
        pass

    messages = [{"role": "system", "content": 
                 "You translate text from third person to first person."}]


    for id, texts in english_texts.items():
        if str(id) in translated_ids:
            print(f"ID: {id} was found to be translated already.\n\n", flush=True)
            continue

        first_person_texts[id] = {}
        for key, text in texts.items():
            if key == "description" or key == "did_you_know":
                try:

                    # Perform conversion from third person to first person using ChatGPT
                    messages.append(
                            {"role": "user", "content": f"Convert the following text to first person:\n\n{text}\n\n"}
                            )
                    chat = openai.ChatCompletion.create(
                        model="gpt-3.5-turbo",
                        messages=messages
                    )
                    reply = chat.choices[0].message.content
                    first_person_texts[id][key] = reply

                    print(f"The ID: {id}\nThe KEY: {key}\nThe TEXT: {text}\nThe REPLY: {reply}\nTries: {tries}\n\n", flush=True)
                    tries += 1
                    if tries > 40:
                        break
                except Exception as e:
                    sleep_timer += 1
                    time.sleep(sleep_timer)
                    print("\n\nError Occurred: \n", e, flush=True)

                    #first_person_texts[id] = first_person_text.choices[0].text.strip()
            else:
                first_person_texts[id][key] = text

        if 'description' in first_person_texts[id] and 'did_you_know' in first_person_texts[id]:
            translated_ids.add(id)
            write_translated_ids(translated_ids, translated_ids_file)
            with open("output.txt", "a") as file:
                file.write(f"Original text from ID: {id}\n")
                file.write(f"Description:\n{english_texts[id]['description']}\n")
                file.write(f"Did You Know:\n{english_texts[id]['did_you_know']}\n\n")
                file.write(f"First-person text from ID: {id}\n")
                file.write(f"Description:\n{first_person_texts[id]['description']}\n")
                file.write(f"Did You Know:\n{first_person_texts[id]['did_you_know']}\n")
                file.write("======================================================================================\n")
            print(f"======================\nCharacter {i} Translated\n======================\n", flush=True)
            i+= 1
            print(f"Sleeper Timer at: {sleep_timer}", flush=True)
            print("+++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++", flush=True)
        if tries > 40:
            break
    print("\n\nDONE WITH EVERYTHING\n\n", flush=True)
    return first_person_texts

'''
Create powerpoint for first person text
'''
def first_person_pptx(first_person_text):
    from origPythonCode import apiRunner, sortDresses, show_error_popup, imageRunner,add_description_subtitle,add_did_you_know_subtitle,add_did_you_know_text,add_image,add_subtitle_highlight, add_title_box,add_description_text,openFile
    try:
        dress_data = apiRunner() # gather all dress data from api
    
        # sort dress data
        sorted_dress_data = sortDresses(dress_data)

        # if there is no image in the local folder then download from web
        if download_imgs.get() == 0:
            if not os.path.exists('./images'):
                os.makedirs('./images')
            if not os.listdir('./images'):
                text_message = "Local folder is empty. Attempting to grab image(s) from web."
                duration_num = 4
                show_error_popup(text_message, duration_num)
                time.sleep(duration_num+1)
                imageRunner(sorted_dress_data)

        # download images from web if download images check box is selected
            # creates directory to save images if one does not exist
        if not os.path.exists('./images'):
            os.makedirs('./images')
        imageRunner(sorted_dress_data) # download images for each dress in list
        prs = Presentation()
        ppt_file_name = "fpp.pptx"
        file_name = "fpp.pptx"
        count = 0
        while os.path.exists(file_name):
            count += 1
            file_name = f"{os.path.splitext(ppt_file_name)[0]}({count}).pptx"

        progress_window = tk.Toplevel(root)
        progress_window.title('Creating First Person Pptx')
        sw = int(progress_window.winfo_screenwidth()/2 - 450/2)
        sh = int(progress_window.winfo_screenheight()/2 - 70/2)
        progress_window.geometry(f'450x70+{sw}+{sh}')
        progress_window.resizable(False, False)
        progress_window.attributes('-disable', True)
        progress_window.focus()

        # progress bar custom style
        pb_style = ttk.Style()
        pb_style.theme_use('clam')
        pb_style.configure('green.Horizontal.TProgressbar', foreground='#1ec000', background='#1ec000')

        # frame to hold progress bar
        pb_frame = tk.Frame(progress_window)
        pb_frame.pack()

        # progress bar
        pb = ttk.Progressbar(pb_frame, length=400, style='green.Horizontal.TProgressbar', mode='determinate', maximum=100, value=0)
        pb.pack(pady=10)

        # label for percent complete
        percent_label = tk.Label(pb_frame, text='Creating Powerpoint...0%')
        percent_label.pack()
        complete = 0

        for id, dress_info in first_person_text.items():
            dress_data = {'id': id}
            print("ID", id, flush=True)
            print("DRESS INFO: ", dress_info, flush=True)
            dress_name = first_person_text[id].get('name', '')
            dress_description = first_person_text[id].get('description', '')
            dress_did_you_know = first_person_text[id].get('did_you_know', '')

            dress_description_len = len(dress_description)



            #--------------------------------Portrait--------------------------------
            prs.slide_width = pptx.util.Inches(7.5) # define slide width
            prs.slide_height = pptx.util.Inches(10.83) # define slide height
            slide_layout = prs.slide_layouts[5] # use slide with only title
            slide_layout2 = prs.slide_layouts[6] # use empty slide

            slide_title = prs.slides.add_slide(slide_layout) 

            add_image(slide_title, dress_data, 0, 1.39)
            add_title_box(slide_title, dress_name, 0, 0.09, 7.5, 0.91) 
            add_subtitle_highlight(slide_title, 3.71, 1.21, 1.83, 0.23) # description - highlight box
            add_description_subtitle(slide_title, 3.63, 0.88, 3.71, 0.37)
            add_description_text(slide_title, dress_description, 3.63, 1.35, 3.71, 7.57)

                # adjust the text height based on text length
            if dress_description_len < 600:
                add_subtitle_highlight(slide_title, 3.72, 5.83, 1.83, 0.23) # did you know - highlight box
                add_did_you_know_subtitle(slide_title, 3.63, 5.42, 3.71, 0.37)
                add_did_you_know_text(slide_title, dress_did_you_know, 3.63, 5.94, 3.71, 0.91)

            elif dress_description_len > 600 and dress_description_len < 1300:
                add_subtitle_highlight(slide_title, 3.72, 7.99, 1.83, 0.23) # did you know - highlight box
                add_did_you_know_subtitle(slide_title, 3.63, 7.58, 3.71, 0.37)
                add_did_you_know_text(slide_title, dress_did_you_know, 3.63, 8.1, 3.71, 0.91)
                
            else:
                add_subtitle_highlight(slide_title, 3.72, 9.32, 1.83, 0.23) # did you know - highlight box
                add_did_you_know_subtitle(slide_title, 3.63, 8.91, 3.71, 0.37)
                add_did_you_know_text(slide_title, dress_did_you_know, 3.63, 9.43, 3.71, 0.91)

            #add_numbering(slide_title, dress_info, id, 0.49, 6.46, 1.33, 0.27, 1.95, 6.46, 1.33, 0.27)
            complete += 1
            #pb['value'] = (complete/len(first_person_text))*100 # calculate percentage of images downloaded
            #percent_label.config(text=f'Creating Book...{int(pb["value"])}%')
            dress_data.clear()
        prs.save(file_name)
        openFile(file_name)
        progress_window.destroy()
    except Exception as e:
        traceback.print_exc()


"""
    Create an HTML package with English and Telugu texts.
    """
def create_html_package(english_texts, telugu_texts):
    page = 1
    html_content = """
    <html>
    <head>
    <style>
        .name { font-weight: bold; text-align: center; }
        .did_you_know { margin-top: 20px; font-size: 18px; }
        .page {
            page-break-after: always;
        }
        @media screen {
            .page {
                border-bottom: 1px solid #ccc;
                padding-bottom: 20px;
                margin-bottom: 20px;
            }
        }
    </style>
    </head>
    <body>
    """
    for id, english_text in english_texts.items():
        telugu_text = telugu_texts.get(id, {'name': '', 'description': '', 'did_you_know': ''})
        html_content += f"<div class='page'><h2>Page No: {page} ABCDid: {id} (English)</h2>"
        html_content += f"<div class='name'>{english_text['name']}</div>"
        html_content += f"<p class='description'>{english_text['description']}</p>"
        html_content += "<div class='did_you_know'>Did you know?</div>"
        html_content += f"<p>{english_text['did_you_know']}</p></div><hr>"
        
        html_content += f"<div class='page'><h2>Page No: {page} ABCDid: {id} (Telugu)</h2>"
        html_content += f"<div class='name'>{telugu_text['name']}</div>"
        html_content += f"<p class='description'>{telugu_text['description']}</p>"
        html_content += "<div class='did_you_know'>Did you know?</div>"
        html_content += f"<p>{telugu_text['did_you_know']}</p></div><hr>"
        page += 1
    html_content += "</body></html>"
    return html_content

def create_html_package_gpt(english_texts, first_person_texts):
    page = 1
    html_content = """
    <html>
    <head>
    <style>
        .name { font-weight: bold; text-align: center; }
        .did_you_know { margin-top: 20px; font-size: 18px; }
        .page {
            page-break-after: always;
        }
        @media screen {
            .page {
                border-bottom: 1px solid #ccc;
                padding-bottom: 20px;
                margin-bottom: 20px;
            }
        }
    </style>
    </head>
    <body>
    """
    for id, english_text in english_texts.items():
        first_person_text = first_person_texts.get(id, {'name': '', 'description': '', 'did_you_know': ''})
        html_content += f"<div class='page'><h2>Page No: {page} ABCD-id: {id} (English)</h2>"
        html_content += f"<div class='name'>{english_text['name']}</div>"
        html_content += f"<div class='description>'>Description</div>"
        html_content += f"<p>{english_text['description']}</p>"
        html_content += "<div class='did_you_know'>Did you know?</div>"
        html_content += f"<p>{english_text['did_you_know']}</p></div><hr>"
        
        html_content += f"<div class='page'><h2>Page No: {page} ABCD-id: {id} (ChatGPT)</h2>"
        html_content += f"<div class='name'>{english_text['name']}</div>"
        html_content += f"<div class='description>'>Description</div>"
        html_content += f"<p>{first_person_text['description']}</p>"
        html_content += "<div class='did_you_know'>Did you know?</div>"
        html_content += f"<p>{first_person_text['did_you_know']}</p></div><hr>"
        page += 1
    html_content += "</body></html>"
    return html_content

"""
    Save the HTML content to a file, appending an incrementing number to the filename
    to avoid overwrites.
    """
def save_html_to_file(english_texts, telugu_texts, base_filename="translation_package"):
    html_filename = f"{base_filename}.html"
    txt_filename = f"{base_filename}.txt"
    counter = 1
    # Check if the file exists and update the filename until it's unique
    while os.path.exists(html_filename) or os.path.exists(txt_filename):
        html_filename = f"{base_filename}_{counter}.html"
        txt_filename = f"{base_filename}_{counter}.txt"
        counter += 1
    
    # Saving HTML content
    with open(html_filename, 'w', encoding='utf-8') as file:
        file.write(create_html_package(english_texts, telugu_texts))  # Assuming create_html_package is defined as before
    
    # Creating a plain text version of the content
    text_content = ""
    for id, english_text in english_texts.items():
        telugu_text = telugu_texts.get(id, '')
        text_content += f"English Text for ID {id}\n{english_text}\n\n"
        text_content += "------------------------------------------------\n\n"
        text_content += f"Telugu Text for ID {id}\n{telugu_text}\n\n"
        text_content += "================================================\n\n"
    
    # Saving text content
    with open(txt_filename, 'w', encoding='utf-8') as file:
        file.write(text_content)

    return html_filename, txt_filename

def save_html_to_file_gpt(english_texts, first_person_texts, base_filename="gpt_adjusted_package"):
    html_filename = f"{base_filename}.html"
    txt_filename = f"{base_filename}.txt"
    counter = 1
    # Check if the file exists and update the filename until it's unique
    while os.path.exists(html_filename) or os.path.exists(txt_filename):
        html_filename = f"{base_filename}_{counter}.html"
        txt_filename = f"{base_filename}_{counter}.txt"
        counter += 1
    
    # Saving HTML content
    with open(html_filename, 'w', encoding='utf-8') as file:
        file.write(create_html_package_gpt(english_texts, first_person_texts))  # Assuming create_html_package is defined as before
    
    # Creating a plain text version of the content
    text_content = ""
    for id, english_text in english_texts.items():
        telugu_text = first_person_texts.get(id, '')
        text_content += f"English Text for ID {id}\n{english_text}\n\n"
        text_content += "------------------------------------------------\n\n"
        text_content += f"ChatGPT Text for ID {id}\n{telugu_text}\n\n"
        text_content += "================================================\n\n"
    
    # Saving text content
    with open(txt_filename, 'w', encoding='utf-8') as file:
        file.write(text_content)

    return html_filename, txt_filename

def generate_translation_package():
    from origPythonCode import apiRunner
    # Get ID numbers from text area
    try:

        dress_data = apiRunner()
        
        # Fetch English text
        english_texts = fetch_english_text(dress_data)
        

        # Translate to Telugu
        telugu_texts = translate_text_to_telugu(english_texts)
        
        # Create HTML package
        html_content = create_html_package(english_texts, telugu_texts)
    
        
        # Save HTML and TXT files
        html_filename, txt_filename = save_html_to_file(english_texts, telugu_texts)
    
        
        # Open the HTML file in a web browser
        webbrowser.open(f'file://{os.path.realpath(html_filename)}')
    
    except FileNotFoundError:
        print(f"File '{e}' not found.")
    except Exception as e:
        print(f'Error: {e}')
    finally:
        translation_package_generate_button.config(state='normal')
        
        # Optionally, show a message that the file has been saved
        print("Translation package has been generated and saved.")



def generate_first_person_package():
    from origPythonCode import apiRunner
    try:

        dress_data = apiRunner() 
        english_texts = fetch_english_text(dress_data)
        first_person_texts = translate_text_to_first_person(english_texts)
        first_person_pptx(first_person_texts)
        create_html_package_gpt(english_texts, first_person_texts)
        html_filename, txt_filename = save_html_to_file_gpt(english_texts, first_person_texts)
        webbrowser.open(f'file://{os.path.realpath(html_filename)}')
    except FileNotFoundError:
        print(f"File '{e}' not found")
    except Exception as e:
        print(f'Error: {e}')
    finally:
        first_person_generate_button.config(state='normal')
    print("Translation package has been generated and saved.", flush=True)
    
def wordSearchOpenAi(english_texts):
    words_for_puzzles = {}
    word_count = int(word_count_var.get())  # Get the current preferred word count

    messages = [
        {"role": "system", "content": f"You will extract {word_count} meaningful words significant to the character's text. Do not reply with anything else. Just list the words."}
    ]

    #sorted_texts = sort_english_texts(english_texts)  # Sort texts based on preference

    for id, text in english_texts.items():
        messages.append(
            {"role": "user", "content": f"Can you extract {word_count} words from this character's text that are meaningful and significant to the character, with one being their name? Please only output the {word_count} words.\n\n{text}\n\n"}
        )
        
        chat_response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=messages
        )
       
        reply = chat_response.choices[0].message.content.strip()
        
        words_for_puzzles[id] = reply
        print(f"This is the Reply for ID {id}: {reply}\n\n")
    
    return words_for_puzzles


#function that splits the words and ids, preparing them for puzzle creation
def wordsearchCreator(words_for_puzzles):
    puzzles = {}
    answer_keys = {}  
    word_lists = {}
    for id, words_string in words_for_puzzles.items():
        word_list = words_string.replace(',', '').replace("'", "").upper().split()
        word_lists[id] = word_list
        grid_size = int(puz_width_var.get())  # Get grid size
        grid = [['-' for _ in range(grid_size)] for _ in range(grid_size)]
        answer_positions = {}

        for word in word_list:
            placed, positions = placeWord(grid, word)
            if placed:
                # Store all positions, not just start and end
                answer_positions[word] = positions

        fillEmptySpots(grid)
        puzzles[id] = grid
        answer_keys[id] = answer_positions
    
    return puzzles, answer_keys, word_lists



#function to randomly place words in the grid
def placeWord(grid, word):
    max_attempts = 100    
    attempts = 0
    placed = False
    positions = []  

    while not placed and attempts < max_attempts:
        wordPlacement = random.randint(0, 3)
        attempts += 1

        if wordPlacement == 0:  # Horizontal
            row = random.randint(0, len(grid) - 1)
            col = random.randint(0, len(grid) - len(word))
            space_available = all(grid[row][col + i] == '-' or grid[row][col + i] == word[i] for i in range(len(word)))
            if space_available:
                for i in range(len(word)):
                    grid[row][col + i] = word[i]
                    positions.append((row, col + i))  
                placed = True

        elif wordPlacement == 1:  # Vertical
            row = random.randint(0, len(grid) - len(word))
            col = random.randint(0, len(grid) - 1)
            space_available = all(grid[row + i][col] == '-' or grid[row + i][col] == word[i] for i in range(len(word)))
            if space_available:
                for i in range(len(word)):
                    grid[row + i][col] = word[i]
                    positions.append((row + i, col))  
                placed = True

        elif wordPlacement == 2:  # Diagonal left to right
            row = random.randint(0, len(grid) - len(word))
            col = random.randint(0, len(grid) - len(word))
            space_available = all(grid[row + i][col + i] == '-' or grid[row + i][col + i] == word[i] for i in range(len(word)))
            if space_available:
                for i in range(len(word)):
                    grid[row + i][col + i] = word[i]
                    positions.append((row + i, col + i)) 
                placed = True

        elif wordPlacement == 3:  # Diagonal right to left
            row = random.randint(0, len(grid) - len(word))
            col = random.randint(len(word) - 1, len(grid) - 1)
            space_available = all(grid[row + i][col - i] == '-' or grid[row + i][col - i] == word[i] for i in range(len(word)))
            if space_available:
                for i in range(len(word)):
                    grid[row + i][col - i] = word[i]
                    positions.append((row + i, col - i))  
                placed = True

    return placed, positions

#function to fill the empty slots after words are placed                 
def fillEmptySpots(grid):
    for row in range(len(grid)):
        for col in range(len(grid[0])):  
            if grid[row][col] == '-':  
                grid[row][col] = random.choice(string.ascii_uppercase)
            
def createWordsearchWordsHtml(puzzles, answer_keys, word_lists):
    page = 1
    html_content = """
    <html>
    <head>
    <title>Word Search Puzzles</title>
    <style>
        body {
            font-family: Arial, sans-serif;
        }
        .page {
            page-break-after: always;
            margin-bottom: 20px;
        }
        table {
            border-collapse: collapse;
            margin: 20px 0;
            position: relative;
            border: 1px solid red; /* Temporary for alignment checking */
        }
        td {
            border: 1px solid #666;
            width: 20px;
            height: 20px;
            text-align: center;
            vertical-align: middle;
            box-sizing: border-box;
            color: black; /* Default text color */
        }
        .answer-text {
            color: blue; /* Change the text color to blue */
            font-weight: bold; /* Make answer text bold */
        }
    </style>
    </head>
    <body>
    """

    # Regular puzzles and word lists
    for id, grid in puzzles.items():
        html_content += f"<div class='page'><h2>Page No: {page} - Puzzle ID: {id}</h2><table>"
        for row in grid:
            html_content += "<tr>"
            for cell in row:
                html_content += f"<td>{cell}</td>"
            html_content += "</tr>"
        html_content += "</table><div><strong>Words:</strong><ul>"
        for word in word_lists[id]:
            html_content += f"<li>{word}</li>"
        html_content += "</ul></div></div>"
        page += 1

    # Answer key puzzles with text color change for answers only
    for id, positions in answer_keys.items():
        grid = puzzles[id]
        html_content += f"<div class='page'><h2>Answer Key Page No: {page} - Puzzle ID: {id}</h2><table>"
        flat_positions = set(sum(positions.values(), []))  # Flatten list of tuples from all words into a set for quick lookup
        for row_idx, row in enumerate(grid):
            html_content += "<tr>"
            for col_idx, cell in enumerate(row):
                # Apply the 'answer-text' class if the cell position is in the flat list of answer positions
                if (row_idx, col_idx) in flat_positions:
                    html_content += f"<td class='answer-text'>{cell}</td>"
                else:
                    html_content += f"<td>{cell}</td>"
            html_content += "</tr>"
        html_content += "</table></div>"
        page += 1

    html_content += "</body></html>"
    return html_content






def save_and_display_html(html_content, base_filename="puzzles_package"):
    html_filename = f"{base_filename}.html"
    counter = 1
    
    # Increment filename if exists to avoid overwriting
    while os.path.exists(html_filename):
        html_filename = f"{base_filename}_{counter}.html"
        counter += 1

    # Save HTML to file
    with open(html_filename, 'w') as file:
        file.write(html_content)
    print(f"HTML content has been saved to {html_filename}.")
    
    # Format the file path for browser compatibility and open it
    try:
        file_url = f"file://{os.path.abspath(html_filename)}"
        webbrowser.open(file_url, new=2)
        print("HTML file has been opened in your web browser.")
    except Exception as e:
        print(f"Failed to open the HTML file in a web browser. Error: {e}")
        
def add_puzzle_table(slide, grid, word_list, title_text, max_width, max_height, answer_positions=None):
    title = slide.shapes.title
    title.text = title_text

    rows, cols = len(grid), len(grid[0])
    grid_origin_x = Inches(1)
    grid_origin_y = Inches(1.5)

    # Calculate the cell width and height based on the max dimensions provided
    cell_width = max_width / cols
    cell_height = max_height / rows

    # Add the grid table
    table = slide.shapes.add_table(rows, cols, grid_origin_x, grid_origin_y, round(cell_width * cols), round(cell_height * rows)).table
    table.first_row = False

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = grid[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)  
            p.alignment = PP_ALIGN.CENTER
            if answer_positions and any((r, c) in positions for positions in answer_positions.values()):
                p.font.color.rgb = RGBColor(255, 0, 0)  

    # Add lines for correct answers using direct line drawing
    if answer_positions:
        for positions in answer_positions.values():
            for i in range(len(positions)-1):
                start_cell = positions[i]
                end_cell = positions[i+1]
                start_x = grid_origin_x + start_cell[1] * cell_width + cell_width / 2
                start_y = grid_origin_y + start_cell[0] * cell_height + cell_height / 2
                end_x = grid_origin_x + end_cell[1] * cell_width + cell_width / 2
                end_y = grid_origin_y + end_cell[0] * cell_height + cell_height / 2

                
                line = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, start_x, start_y, end_x - start_x, end_y - start_y)
                line.line.width = Pt(2)
                line.line.color.rgb = RGBColor(255, 0, 0)  

    # Add word list to the side
    textbox = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(2), Inches(4))
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "Words to find:\n" + "\n".join(word_list)
    p.font.bold = True
    p.font.size = Pt(14)




def make_powerpoint(puzzles, answer_keys, word_lists):
    prs = Presentation()
    max_width = Inches(6)  # Set the maximum width for the grid
    max_height = Inches(4.5)  # Set the maximum height for the grid

    for id, grid in puzzles.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        add_puzzle_table(slide, grid, word_lists[id], f"Puzzle ID: {id}", max_width, max_height)

    for id, positions in answer_keys.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        add_puzzle_table(slide, puzzles[id], word_lists[id], f"Answer Key ID: {id}", max_width, max_height, answer_positions=positions)

    return prs






def save_powerpoint(prs, base_filename="puzzles_package"):
    filename = f"{base_filename}.pptx"
    counter = 1
    
    # Increment filename if exists to avoid overwriting
    while os.path.exists(filename):
        filename = f"{base_filename}_{counter}.pptx"
        counter += 1

    # Save PowerPoint to file
    prs.save(filename)
    print(f"PowerPoint content has been saved to {filename}.")
    
def generate_word_search_package():
    from origPythonCode import apiRunner
    
    dress_data = apiRunner()

    
    english_texts = fetch_english_text(dress_data)

   
    puzzleWords = wordSearchOpenAi(english_texts)

    
    puzzles, answer_keys, word_lists = wordsearchCreator(puzzleWords)

    
    html_puzzles = createWordsearchWordsHtml(puzzles, answer_keys, word_lists)

   
    save_and_display_html(html_puzzles)
    print("HTML word puzzles have been generated and saved.")

   
    prs = make_powerpoint(puzzles, answer_keys, word_lists)

    
    save_powerpoint(prs, base_filename="word_puzzles_package")
    print("PowerPoint word puzzles have been generated and saved.")






'''
    ##DOB Analayser Methods
'''
def parse_dates(dates):
    if pd.isna(dates) or dates.strip() in ["TBD", "Unknown", "Not Applicable", ""]:
        # print(f"Skipping date parsing due to non-date info: {dates}")
        return None, None

    # Normalize and clean the date string
    dates = dates.replace('(', '').replace(')', '').replace('circa', '').replace('c.', '').replace('around', '').strip()
    dates = dates.replace(' – ', '-').replace(' to ', '-').replace('–', '-').replace('/', '-')

    # Split and strip the date parts
    try:
        split_dates = [date.strip() for date in dates.split('-') if date.strip()]
        birth_date = parse_single_date(split_dates[0])
        death_date = parse_single_date(split_dates[1]) if len(split_dates) > 1 else None
    except Exception as e:
        return None, None

    return birth_date, death_date

def parse_single_date(date_str):
    try:
        return datetime.datetime.strptime(date_str, '%d %B %Y').date()
    except ValueError:
        try:
            return datetime.datetime.strptime(date_str, '%B %d, %Y').date()
        except ValueError:
            return None

def calculate_life_span(birth_date, death_date):
    if not birth_date:
        return "Birth date unknown"
    if not death_date:
        return "Still living"  # Or adjust according to your preference
    try:
        life_span = death_date.year - birth_date.year
        life_span -= ((death_date.month, death_date.day) < (birth_date.month, birth_date.day))
        return life_span
    except Exception as e:
        return ""

def load_data(file_path):
    data = pd.read_excel(file_path)
    data['Parsed Dates'] = data['Dates'].apply(parse_dates)
    data[['Date of Birth', 'Date of Death']] = pd.DataFrame(data['Parsed Dates'].tolist(), index=data.index)
    data['Life Span'] = data.apply(lambda row: calculate_life_span(row['Date of Birth'], row['Date of Death']), axis=1)
    return data[['name', 'abcd_id', 'Date of Birth', 'Date of Death', 'Life Span']]

def upload_file():
    file_path = filedialog.askopenfilename()
    if file_path:
        global data_frame
        data_frame = load_data(file_path)
        print("File uploaded and processed successfully.")

def generate_html():
    if data_frame.empty:
        # print("No data to display.")
        return
    html_content = '<html><head><title>Excel Data</title>'
    html_content += '<script src="https://code.jquery.com/jquery-3.5.1.js"></script>'
    html_content += '<script src="https://cdn.datatables.net/1.10.21/js/jquery.dataTables.min.js"></script>'
    html_content += '<link rel="stylesheet" href="https://cdn.datatables.net/1.10.21/css/jquery.dataTables.min.css">'
    html_content += '<script>$(document).ready(function(){$("#data_table").DataTable();});</script></head>'
    html_content += '<body><table id="data_table" class="display" style="width:100%"><thead><tr>'
    html_content += '<th>ID</th><th>Name</th><th>Date of Birth</th><th>Date of Death</th><th>Life Span</th></tr></thead><tbody>'
    for _, row in data_frame.iterrows():
        html_content += f'<tr><td>{row["abcd_id"]}</td><td>{row["name"]}</td><td>{row.get("Date of Birth", "")}</td><td>{row.get("Date of Death", "")}</td><td>{row.get("Life Span", "")}</td></tr>'
    html_content += '</tbody></table></body></html>'
    with tempfile.NamedTemporaryFile('w', delete=False, suffix='.html') as f:
        f.write(html_content)
        webbrowser.open('file://' + f.name)



'''
Play Audio
'''
def playAudio():
    try:
        engine = pyttsx3.init()
        text = text_field_Description.get("1.0", tk.END)
        engine.say(text)
        engine.runAndWait()
    except Exception as e:
       messagebox.showerror("Play Audio Error", str(e))

'''
Save Audio
'''
def saveAudio():
    try:
        engine = pyttsx3.init()
        text = text_field_Description.get("1.0", tk.END)
        #id = text_field_ID.get()
        engine.save_to_file(text, f"{text_field_ID.get()}.mp3")
        engine.runAndWait()
        messagebox.showinfo("Save Audio", "Audio saved as id.mp3")
    except Exception as e:
        messagebox.showerror("Save Audio Error", str(e))

def SaveAllAudios():
    update_dress_list = []
    # get dress numbers from text field
    get_text_field = text_field.get("1.0", "end-1c").split(',')

    # add to list
    for number in get_text_field:
        if (number.strip().isnumeric()):
            update_dress_list.append(int(number.strip()))
    for num in update_dress_list:
        print(num)
        fetchTextAndSaveAudio(num)

def fetchTextAndSaveAudio(id):
    headers = {
        'Accept': '*/*',
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/98.0.4758.102 Safari/537.36'
    }
    try:
        # Replace with the actual API endpoint you're supposed to hit
        url = f"https://abcd2.projectabcd.com/api/getinfo.php?id={id}"
        response = requests.get(url, headers=headers)
        if response.ok:
            data = response.json()
            description = data['data']['description']

            # Convert the description to audio
            tts = gTTS(description, lang='en')
            # Save the audio file named as id.mp3
            audio_file_path = f"{id}.mp3"
            tts.save(audio_file_path)
        else:
            print(f"Failed to fetch data for ID {id}: {response.status_code}, {response.reason}")
    except Exception as e:
        print(f"Error fetching data for ID {id}: {e}")

'''
Get Text
'''
def fetchText(id):
    headers = {
        'Accept': '*/*',  # Use wildcard or specific type based on API requirement
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/98.0.4758.102 Safari/537.36'
    }
    try:
        url = f"https://abcd2.projectabcd.com/api/getinfo.php?id={id}"
        response = requests.get(url, headers=headers)
        # print("Request Headers:", response.request.headers)
        if response.ok:
            # print("API Response:", response.text)
            data = response.json()
            # description = data.get("description", "No description found.")
            description = data['data']['description']
            text_field_Description.delete(1.0, tk.END)  # This clears all text from the widget
            root.after(0, lambda: text_field_Description.insert(tk.END, description))
        else:
            print("Failed to fetch data:", response.status_code, response.reason)
            messagebox.showerror("Error", f"Failed to fetch data from the server: {response.reason}, Status Code: {response.status_code}")
    except Exception as e:
        traceback.print_exc()
        print("Error:", str(e))
        messagebox.showerror("Error", str(e))

'''
Spins up new thread to run getText function
'''
def getTextThread():
    dress_id = text_field_ID.get()
    if dress_id:
        thread = threading.Thread(target=fetchText, args=(dress_id,))
        thread.start()
    else:
        messagebox.showerror("Error", "Please enter a dress ID.")
        
######################################################################THREAD SECTION########################################################################################

'''
Spins up new thread to run generateBook
'''
def startGenerateBookThread():
    from origPythonCode import generateBook
    book_gen_generate_button.config(state="disabled")
    generate_thread = threading.Thread(target=generateBook)
    generate_thread.start()

'''
Spins up new thread to run diffReport function
'''
def startDiffReportThread():
    from origPythonCode import diffReport
    diff_report_button.config(state='disabled')
    diff_report_thread = threading.Thread(target=diffReport)
    diff_report_thread.start()

'''
Spins up new thread to run wordAnalysis function
'''
def startWordAnalysisThread():
    from origPythonCode import wordAnalysis
    word_analysis_button.config(state='disabled')
    word_analysis_thread = threading.Thread(target=wordAnalysis)
    word_analysis_thread.start()

'''
Spins up new thread to run googleImage function
'''
def startGoogleImageThread():
    from origPythonCode import googleImage
    google_image_search_button.config(state='disabled')
    google_image_search_thread = threading.Thread(target=googleImage)
    google_image_search_thread.start()

'''
Spins up new thread to run generateWikiLink function
'''
def startGenerateWikiLinkThread():
    from origPythonCode import generateWikiLink
    wiki_link_gen_button.config(state='disabled')
    wiki_link_thread = threading.Thread(target=generateWikiLink)
    wiki_link_thread.start()

'''
Spins up new thread to run generatePairs function
'''
def startGeneratePairsThread():
    who_are_my_pairs_gen_button.config(state='disabled')
    who_are_my_pairs_thread = threading.Thread(target=generatePairs)
    who_are_my_pairs_thread.start()

'''
Spins up new thread to run us/uk_spellings function
'''
def startUS_UK_SpellingsThread():
    us_uk_spellings_button.config(state="disabled")
    us_uk_spellings_thread = threading.Thread(target=generate_us_uk_spellings)
    us_uk_spellings_thread.start()

'''
Spins up new thread to run ids of us/uk_spellings function
'''
def startIDs_Of_US_UK_SpellingsThread():
    ids_of_us_uk_spellings_button.config(state="disabled")
    ids_of_us_uk_spellings_thread = threading.Thread(target=generate_IDs_Of_us_uk_spellings)
    ids_of_us_uk_spellings_thread.start()


'''
Spins up new thread to run translatepackage function
'''
def startTranslationPackageThread():
    translation_package_generate_button.config(state="disabled")
    translate_package_thread = threading.Thread(target=generate_translation_package)
    translate_package_thread.start()

'''
Spins up new thread to run translate_to_first person function
'''
def startFirstPersonThread():
    first_person_generate_button.config(state="disabled")
    first_person_thread = threading.Thread(target=generate_first_person_package)
    first_person_thread.start()

'''
Spins up new thread to run playAudio function
'''
def playAudioThread():
    get_audio_button.config(state='disabled')
    get_audio_thread = threading.Thread(target=playAudio)
    get_audio_thread.start()

'''
Spins up new thread to run saveAudio function
'''
def saveAudioThread():
    get_audio_button.config(state='disabled')
    get_audio_thread = threading.Thread(target=saveAudio)
    get_audio_thread.start()
'''
Spins up new thread to run word search function
'''
def startWordPuzzleThread():
    word_puzzle_generate_button.config(state="disabled")
    word_search_thread = threading.Thread(target=generate_word_search_package)
    word_search_thread.start()


'''
Launch help site when user clicks Help button
'''
def launchHelpSite():
    # create help site
    with open('help.html', 'w') as file:
        file.write('<!DOCTYPE html>\n<html>\n<head>\n\t<meta charset="utf8">\n\t<title>abcd Help</title>\n</head>\n<body>\n\t\t<h1 style="text-align: center;">Welcome to the help site</h1>\n</body>\n</html>\n')
    
    # open help site
    webbrowser.open('help.html')

'''
Raise selected frame to the top
'''
def raiseFrame(frame):
    if frame == 'main_frame':
        main_frame.tkraise()
        root.title("Project ABCD Admin Panel")
    elif frame == 'book_gen_frame':
        book_gen_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Book Generation")
    elif frame == 'diff_report_frame':
        diff_report_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Difference Report")
    elif frame == 'word_analysis_frame':
        word_analysis_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Word Analysis")
    elif frame == 'google_image_frame':
        google_image_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Google Image")
    elif frame == 'wiki_link_frame':
        wiki_link_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Wiki Link")
    elif frame == 'who_are_my_pairs_frame':
        who_are_my_pairs_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Who Are My Pairs")
    elif frame == 'translation_package_frame':  
        translation_package_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Translation Package")
    elif frame == 'first_person_frame':
        first_person_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
    elif frame == 'us_uk_spellings_frame':
        us_uk_spellings_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD US/UK Spellings")
    elif frame == 'ids_of_us_uk_spellings_frame':
        ids_of_us_uk_spellings_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD IDs_Of_US/UK Spellings")
    elif frame == 'Get_audio_frame':
        Get_audio_frame.tkraise()
        text_field_label_ID_Address.tkraise()
        # text_field_label_ID_Address.tkraise()
        text_field_label_Description.tkraise()
        text_field_Description.tkraise()
        text_field_ID.tkraise()
        play_audio_button.tkraise()
        save_audio_button.tkraise()
        upload_audio_button.tkraise()
        get_text_button.tkraise()
        root.title("Project ABCD Get Audio")
    elif frame == "all_audio_frame":
        all_audio_frame.tkraise()
        all_audio_button.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Word Analysis")
    elif frame == "DOB_Analyzer_frame":
        DOB_Analyzer_frame.tkraise()
        word_analysis_button.tkraise()
        word_analysis_back_button.tkraise()
        UPLOAD_FILE_button.tkraise()
        # UPLOAD_FILE_button.tkraise()
        root.title("Project ABCD Word Analysis")
    elif frame == 'word_puzzle_frame':
        word_puzzle_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Word Search Puzzles")


#--------------------------------Main Frame-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------------------------------------------------
# configure grid to fill extra space and center
tk.Grid.rowconfigure(root, 0, weight=1)
tk.Grid.columnconfigure(root, 0, weight=1)

# main frame
main_frame = tk.Frame(root, width=1000, height=1000)
main_frame.pack_propagate(False)
main_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Main Title-----------------------------------------------------------------------------------------------
# Create a label widget
title_label = tk.Label(main_frame, text="Project ABCD\nMain Menu",  font=('Arial', 20))
title_label.pack(pady=100)

#--------------------------------Main Buttons-----------------------------------------------------------------------------------------------
# Create buttons widget
## Button settings
main_button_frame = tk.Frame(main_frame)
main_button_frame.place(relx=.5, rely=.35, anchor='center')
button_width = 20
button_height = 3
button_bgd_color = "#007FFF"
button_font_color = "#ffffff"

## Generate Book: Gets selected dress from API and import into ppt
generate_book_button = tk.Button(main_button_frame, text="Generate Book", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('book_gen_frame'))
generate_book_button.pack(side="left", padx=50)

## Diff Report: Create a SQL file of dresses that got changed from excel sheet byt comparing to API
diff_report_button = tk.Button(main_button_frame, text="Difference Report", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('diff_report_frame'))
diff_report_button.pack(side="left", padx=50, anchor='center')

## Generate Book: Get selected dress that user input & put into a table (ID, Name, Description Count, DYK Count, Total Nouns Count, Total Adjectives Count)
word_analysis_report_button = tk.Button(main_button_frame, text="Word Analysis Report", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('word_analysis_frame'))
word_analysis_report_button.pack(side="left", padx=50)

main_button_frame2 = tk.Frame(main_frame)
main_button_frame2.place(relx=.5, rely=.50, anchor='center')

## Google Images: Create an Excel file with 3 image links to the selected dresses
google_image_button = tk.Button(main_button_frame2, text="Google Image", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('google_image_frame'))
google_image_button.pack(side="left", padx=50)

## Wiki Link: [FILL IN THE ACTION HERE]
wiki_link_button = tk.Button(main_button_frame2, text="Wiki Link", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('wiki_link_frame'))
wiki_link_button.pack(side="left", padx=50)

## My Pairs: Shows pairs when they are searched
who_are_my_pairs_button = tk.Button(main_button_frame2, text="Who Are My Pairs?", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('who_are_my_pairs_frame'))
who_are_my_pairs_button.pack(side="left", padx=50)

main_button_frame3 = tk.Frame(main_frame)
main_button_frame3.place(relx=.5, rely=.65, anchor='center')

## US Spelling: us_uk spellings
us_uk_spellings_button = tk.Button(main_button_frame3, text="US/UK Spellings", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('us_uk_spellings_frame'))
us_uk_spellings_button.pack(side="left", padx=50)

## IDS OF US Spelling: ids of us_uk spellings
ids_of_us_uk_spellings_button = tk.Button(main_button_frame3, text="IDS OF US/UK Spellings", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('ids_of_us_uk_spellings_frame'))
ids_of_us_uk_spellings_button.pack(side="left", padx=50)

## Generate Book: fetches English text from Api, uses google translate to generate "telugu" text, then creates HTML package with english text on page, and "telugu" text on another.
translation_package_button = tk.Button(main_button_frame3, text="Translation Package", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('translation_package_frame'))
translation_package_button.pack(side="left", padx=50)

main_button_frame4 = tk.Frame(main_frame)
main_button_frame4.place(relx=.5, rely=.80, anchor='center')

## First Person: fetches text from api, uses ChatGPT to reword the description and did you know text to first person
first_person_button = tk.Button(main_button_frame4, text="First Person Conversion", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('first_person_frame'))
first_person_button.pack(side="left", padx=50)

## Word Puzzle: generates and creates crossword puzzles based of words in character descriptions
word_puzzle_button = tk.Button(main_button_frame4, text="Word Puzzle Creator", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('word_puzzle_frame'))
word_puzzle_button.pack(side="left", padx=50)

main_button_frame5 = tk.Frame(main_frame)
main_button_frame5.place(relx =.5, rely=.95, anchor='center')

##Get Audio Frame
get_audio_button = tk.Button(main_button_frame5, text="Get Audio", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('Get_audio_frame'))
get_audio_button.pack(side="left", padx=50)

## get all Audio Frame
get_all_audio_button = tk.Button(main_button_frame5, text="Get All Audio", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('all_audio_frame'))
get_all_audio_button.pack(side="left", padx=50)

## get DOB Analyzer Frame
DOB_Analyzer_button = tk.Button(main_button_frame4, text="DOB Analyzer", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('DOB_Analyzer_frame'))
DOB_Analyzer_button.pack(side="left", padx=50)

#--------------------------------Book Gen Frame---------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------------------------------------------------------
# book gen frame
book_gen_frame = tk.Frame(root, width=1000, height=600)
book_gen_frame.pack_propagate(False)
book_gen_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Text Field-----------------------------------------------------------------------------------------------
# text field label
text_field_label = tk.Label(root, text="Dress Numbers:", font=LABEL_FONT)
text_field_label.place(x=25, y=72.5)

# text field
text_field = tk.Text(root)
text_field.place(x=175, y=10, relwidth=.8, height=135)

# text field initialization
try:
    with open('slide_numbers.txt', 'r') as file:
        slide_number_content = file.readline().strip()
        text_field.insert("1.0", slide_number_content)
except FileNotFoundError:
    print(FileNotFoundError)

#--------------------------------Layout Radio Buttons-------------------------------------------------------------------------------------
# layout variable
layout = tk.IntVar()
layout.set(4)

# layout frame
layout_frame = tk.Frame(book_gen_frame)
# layout radio buttons
layout_radio4 = tk.Radiobutton(layout_frame, text="Picture on Left - Text on right - Single Page - Portrait Mode",
                                font=MAIN_FONT, variable=layout, value=4)
layout_radio1 = tk.Radiobutton(layout_frame, text="Picture on Left Page - Text on Right Page - Two Page Mode - Portrait Mode",
                                font=MAIN_FONT, variable=layout, value=1)
layout_radio2 = tk.Radiobutton(layout_frame, text="Picture on Right - Text on Left - Single Page - Landscape Mode",
                                font=MAIN_FONT, variable=layout, value=2)
layout_radio3 = tk.Radiobutton(layout_frame, text="Picture on Left - Text on Right - Single Page - Landscape Mode",
                                font=MAIN_FONT, variable=layout, value=3)
# pack radio buttons into layout frame
layout_radio4.pack(anchor="nw")
layout_radio1.pack(anchor="nw")
layout_radio2.pack(anchor="nw")
layout_radio3.pack(anchor="nw")
# place layout frame on main frame
layout_frame.place(x=175, y=150, width=800)
# layout radio buttons label
layout_label = tk.Label(book_gen_frame, text="Layout:", font=LABEL_FONT)
layout_label.place(x=25, y=170)

# separator line
separator1 = ttk.Separator(book_gen_frame)
separator1.place(x=175, y=150, relwidth=.8)

#--------------------------------Sort Radio Buttons---------------------------------------------------------------------------------------
# sort variable
sort_order = tk.IntVar()
sort_order.set(1)

# sort frame
sort_frame = tk.Frame(book_gen_frame)
# sort radio buttons
sort_radio1 = tk.Radiobutton(sort_frame, text="By Name", font=MAIN_FONT, variable=sort_order, value=1)
sort_radio2 = tk.Radiobutton(sort_frame, text="By ID", font=MAIN_FONT, variable=sort_order, value=2)
sort_radio3 = tk.Radiobutton(sort_frame, text="By Input Order", font=MAIN_FONT, variable=sort_order, value=3)
# pack radio buttons into sort frame
sort_radio1.pack(side="left")
sort_radio2.pack(side="left")
sort_radio3.pack(side="left")
# place sort frame on main frame
sort_frame.place(x=175, y=265, width=800)

# sort radio buttons label
sort_label = tk.Label(book_gen_frame, text="Sort Order:", font=LABEL_FONT)
sort_label.place(x=25, y=265)

# separator line
separator2 = ttk.Separator(book_gen_frame)
separator2.place(x=175, y=265, relwidth=.8)


#--------------------------------Preferences----------------------------------------------------------------------------------------------
# Initialize variables for storing values
text_size_var = tk.StringVar()
title_size_var = tk.StringVar()
subtitle_size_var = tk.StringVar()
text_font_var = tk.StringVar()
title_font_var = tk.StringVar()
subtitle_font_var = tk.StringVar()
pic_width_var = tk.StringVar()
pic_height_var = tk.StringVar()

# Set initial values for the Entry fields
text_size_var.set(preferences["TEXT_SIZE"])
title_size_var.set(preferences["TITLE_SIZE"])
subtitle_size_var.set(preferences["SUBTITLE_SIZE"])
text_font_var.set(preferences["TEXT_FONT"])
title_font_var.set(preferences["TITLE_FONT"])
subtitle_font_var.set(preferences["SUBTITLE_FONT"])
pic_width_var.set(preferences["PIC_WIDTH"])
pic_height_var.set(preferences["PIC_HEIGHT"])

# preferences frame
preferences_frame = tk.Frame(book_gen_frame)
# preferences labels and entry fields
text_size_label = tk.Label(preferences_frame, text="Text Size:", font=MAIN_FONT)
text_size = tk.Entry(preferences_frame, width=3, textvariable=text_size_var, state="disabled", font=MAIN_FONT)

title_size_label = tk.Label(preferences_frame, text="Title Size:", font=MAIN_FONT)
title_size = tk.Entry(preferences_frame, width=3, textvariable=title_size_var, state="disabled",  font=MAIN_FONT)

subtitle_size_label = tk.Label(preferences_frame, text="Subtitle Size:", font=MAIN_FONT)
subtitle_size = tk.Entry(preferences_frame, width=3, textvariable=subtitle_size_var, state="disabled",  font=MAIN_FONT)

text_font_label = tk.Label(preferences_frame, text="Text Font:", font=MAIN_FONT)
text_font = tk.Entry(preferences_frame, width=25, textvariable=text_font_var, state="disabled",  font=MAIN_FONT)

title_font_label = tk.Label(preferences_frame, text="Title Font:", font=MAIN_FONT)
title_font = tk.Entry(preferences_frame, width=25, textvariable=title_font_var, state="disabled",  font=MAIN_FONT)

subtitle_font_label = tk.Label(preferences_frame, text="Subitle Font:", font=MAIN_FONT)
subtitle_font = tk.Entry(preferences_frame, width=25, textvariable=subtitle_font_var, state="disabled",  font=MAIN_FONT)

pic_width_label = tk.Label(preferences_frame, text="Pic Width:", font=MAIN_FONT)
pic_width = tk.Entry(preferences_frame, width=6, textvariable=pic_width_var, state="disabled",  font=MAIN_FONT)

pic_height_label = tk.Label(preferences_frame, text="Pic Height:", font=MAIN_FONT)
pic_height = tk.Entry(preferences_frame, width=6, textvariable=pic_height_var, state="disabled",  font=MAIN_FONT)

# grid preference labels and entry fields into preferences frame
# column 1 + 2
text_size_label.grid(row=1, column=1, pady=10)
text_size.grid(row=1, column=2)

title_size_label.grid(row=2, column=1, pady=10)
title_size.grid(row=2, column=2)

subtitle_size_label.grid(row=3, column=1, pady=10, padx=15)
subtitle_size.grid(row=3, column=2)

# column 3 + 4
text_font_label.grid(row=1, column=3)
text_font.grid(row=1, column=4)

title_font_label.grid(row=2, column=3)
title_font.grid(row=2, column=4)

subtitle_font_label.grid(row=3, column=3, padx=15)
subtitle_font.grid(row=3, column=4)

# column 5 + 6
pic_width_label.grid(row=1, column=5)
pic_width.grid(row=1, column=6)

pic_height_label.grid(row=2, column=5, padx=15)
pic_height.grid(row=2, column=6)

# place preferences frame on main frame
preferences_frame.place(x=175, y=295, width=800)

# preferences label
preferences_label = tk.Label(book_gen_frame, text="Preferences:", font=LABEL_FONT)
preferences_label.place(x=25, y=350)

# separator line
separator3 = ttk.Separator(book_gen_frame)
separator3.place(x=175, y=295, relwidth=.8)

#--------------------------------Numbering Radio Buttons--------------------------------------------------------------------------------------
# numbering variable
numbering = tk.IntVar()
numbering.set(1)

# numbering frame
numbering_frame = tk.Frame(book_gen_frame)
# numbering radio buttons
numbering_radio1 = tk.Radiobutton(numbering_frame, text="Show both Page Number and Dress ID", font=MAIN_FONT, variable=numbering, value=1)
numbering_radio2 = tk.Radiobutton(numbering_frame, text="Show Page Number", font=MAIN_FONT, variable=numbering, value=2)
numbering_radio3 = tk.Radiobutton(numbering_frame, text="Show Dress ID", font=MAIN_FONT, variable=numbering, value=3)
# pack numbering buttons into sort frame
numbering_radio1.pack(side="left")
numbering_radio2.pack(side="left")
numbering_radio3.pack(side="left")
# place numbering frame on main frame
numbering_frame.place(x=175, y=445, width=800)

# numbering radio buttons label
numbering_label = tk.Label(book_gen_frame, text="Numbering:", font=LABEL_FONT)
numbering_label.place(x=25, y=445)

# separator line
separator4 = ttk.Separator(book_gen_frame)
separator4.place(x=175, y=445, relwidth=.8)

#--------------------------------Translate and Image Check Buttons--------------------------------------------------------------------------------------
# translate check variable
translate = tk.IntVar()
translate.set(0)
# language options variable
language = tk.StringVar()
language.set(LANGUAGES[0])
# Download image variable
download_imgs = tk.IntVar()
download_imgs.set(0)
# Generate book from local Excel sheet
gen_local = tk.IntVar()
gen_local.set(0)

# translate frame
check_button_frame = tk.Frame(book_gen_frame)
# translate check button
translate_checkbutton = tk.Checkbutton(check_button_frame, text="Translate to:", font=MAIN_FONT, variable=translate, onvalue=1, offvalue=0)
# language options
language_options = tk.OptionMenu(check_button_frame, language, *LANGUAGES)
# download images
download_images = tk.Checkbutton(check_button_frame, text="Download Images", font=MAIN_FONT, variable=download_imgs, onvalue=1, offvalue=0, command=lambda: gen_local.set(0))
# generate book from local Excel sheet
generate_from_local = tk.Checkbutton(check_button_frame, text="Generate from Local", font=MAIN_FONT, variable=gen_local, onvalue=1, offvalue=0, command=lambda: download_imgs.set(0))

# pack translate options into translate frame
translate_checkbutton.pack(side="left")
language_options.pack(side="left")
download_images.pack(side="left")
generate_from_local.pack(side="left")
# place translate frame on main frame
check_button_frame.place(x=175, y=495)

# separator line
separator5 = ttk.Separator(book_gen_frame)
separator5.place(x=175, y=495, relwidth=.8)

#--------------------------------Book Gen Buttons--------------------------------------------------------------------------------------
# button frame
book_gen_button_frame = tk.Frame(book_gen_frame)
# generate button
book_gen_generate_button = tk.Button(book_gen_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startGenerateBookThread)
# help button
book_gen_help_button = tk.Button(book_gen_button_frame, text="Help", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=launchHelpSite)
# upload button
book_gen_back_button = tk.Button(book_gen_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
# pack buttons into button frame
book_gen_generate_button.pack(side="left", padx=35)
book_gen_help_button.pack(side="left")
book_gen_back_button.pack(side="left", padx=30)

# place button frame on main frame
book_gen_button_frame.pack(side="bottom", pady=10)


#--------------------------------Diff Report Frame-----------------------------------------------------------------------------------------------
#------------------------------------------------------------------------------------------------------------------------------------------------
diff_report_frame = tk.Frame(root, width=1000, height=600)
diff_report_frame.pack_propagate(False)
diff_report_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Diff Report Buttons-----------------------------------------------------------------------------------------------
# button frame
diff_report_button_frame = tk.Frame(diff_report_frame)
# difference report button
diff_report_button = tk.Button(diff_report_button_frame, text="Diff Report", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startDiffReportThread)
# back button
diff_back_button = tk.Button(diff_report_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
# pack buttons into button frame
diff_report_button.pack(side="left", padx=35)
diff_back_button.pack(side="left", padx=30)

# place button frame on diff report frame
diff_report_button_frame.pack(side="bottom", pady=10)


#--------------------------------Word Analysis Frame-----------------------------------------------------------------------------------------------
#--------------------------------------------------------------------------------------------------------------------------------------------------
word_analysis_frame = tk.Frame(root, width=1000, height=600)
word_analysis_frame.pack_propagate(False)
word_analysis_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Word Analysis Buttons-----------------------------------------------------------------------------------------------
# button frame
word_analysis_button_frame = tk.Frame(word_analysis_frame)
# word analysis button
word_analysis_button = tk.Button(word_analysis_button_frame, text="Word Analysis", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startWordAnalysisThread)
# back button
word_analysis_back_button = tk.Button(word_analysis_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
# pack buttons into button frame
word_analysis_button.pack(side="left", padx=35)
word_analysis_back_button.pack(side="left", padx=30)

# place button frame on word analysis frame
word_analysis_button_frame.pack(side="bottom", pady=10)

#--------------------------------Google Image Frame-----------------------------------------------------------------------------------------------
#--------------------------------------------------------------------------------------------------------------------------------------------------
google_image_frame = tk.Frame(root, width=1000, height=600)
google_image_frame.pack_propagate(False)
google_image_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Google Image Buttons-----------------------------------------------------------------------------------------------
# Download google image variable
download_google_imgs = tk.IntVar()
download_google_imgs.set(0)

# button frame
google_image_button_frame = tk.Frame(google_image_frame)
# checkbox frame
download_google_image_checkbutton = tk.Frame(google_image_frame)
# download images checkbox
download_google_images = tk.Checkbutton(download_google_image_checkbutton, text="Download Images", font=MAIN_FONT, variable=download_google_imgs, onvalue=1, offvalue=0, command=lambda: gen_local.set(0))
# google search button
google_image_search_button = tk.Button(google_image_button_frame, text="Google Search", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startGoogleImageThread)
# back button
google_image_back_button = tk.Button(google_image_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack buttons into button frame
download_google_image_checkbutton.pack(side="left")
download_google_images.pack(side="left")
google_image_search_button.pack(side="left", padx=35)
google_image_back_button.pack(side="left", padx=30)
# place google image frame on main frame
google_image_button_frame.place(x=175, y=495)
download_google_image_checkbutton.place(x=170, y=150)

# place button frame on word analysis frame
google_image_button_frame.pack(side="bottom", pady=10)

#--------------------------------Wiki Link Frame-----------------------------------------------------------------------------------------------
#------------------------------------------------------------------------------------------------------------------------------------------------
wiki_link_frame = tk.Frame(root, width=1000, height=600)
wiki_link_frame.pack_propagate(False)
wiki_link_frame.grid(row=0, column=0, sticky='news')
 
#--------------------------------Wiki Link Buttons--------------------------------------------------------------------------------------
# button frame
wiki_link_gen_button_frame = tk.Frame(wiki_link_frame)

wiki_link_gen_button = tk.Button(wiki_link_gen_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startGenerateWikiLinkThread)
wiki_link_back_button = tk.Button(wiki_link_gen_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack buttons into button frame
wiki_link_gen_button.pack(side="left", padx=35)
wiki_link_back_button.pack(side="left", padx=30)

# place button frame on word analysis frame
wiki_link_gen_button_frame.pack(side="bottom", pady=10)

#--------------------------------Who Are My Pairs Frame-----------------------------------------------------------------------------------------------
#------------------------------------------------------------------------------------------------------------------------------------------------
who_are_my_pairs_frame = tk.Frame(root, width=1000, height=600)
who_are_my_pairs_frame.pack_propagate(False)
who_are_my_pairs_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Who Are My Pairs Buttons-----------------------------------------------------------------------------------
#button frame
who_are_my_pairs_gen_button_frame = tk.Frame(who_are_my_pairs_frame)

who_are_my_pairs_gen_button = tk.Button(who_are_my_pairs_gen_button_frame, text="Generate Pairs", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startGeneratePairsThread)
who_are_my_pairs_back_button = tk.Button(who_are_my_pairs_gen_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack button into button frame
who_are_my_pairs_gen_button.pack(side="left", padx=35)
who_are_my_pairs_back_button.pack(side="left", padx=30)

# place button frame on word analysis frame?
who_are_my_pairs_gen_button_frame.pack(side="bottom", pady=10)

#--------------------------------Translation Package Frame-----------------------------------------------------------------------------------------------
#--------------------------------------------------------------------------------------------------------------------------------------------------
translation_package_frame = tk.Frame(root, width=1000, height=600)
translation_package_frame.pack_propagate(False)
translation_package_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Translation Package Buttons-----------------------------------------------------------------------------------------------
# generate button frame
translation_package_button_frame = tk.Frame(translation_package_frame)
# generate button
translation_package_generate_button = tk.Button(translation_package_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startTranslationPackageThread)
# help button
translation_package_help_button = tk.Button(translation_package_button_frame, text="Help", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=launchHelpSite)
# upload button
translation_package_back_button = tk.Button(translation_package_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
# pack buttons into button frame
translation_package_generate_button.pack(side="left", padx=35)
translation_package_help_button.pack(side="left")
translation_package_back_button.pack(side="left", padx=30)

# place button frame on word analysis frame
translation_package_button_frame.pack(side="bottom", pady=10)

#-------------------------------First Person Frame-------------------------------------------------------------------------------------------
first_person_frame = tk.Frame(root, width=1000, height=600)
first_person_frame.pack_propagate(False)
first_person_frame.grid(row=0, column=0, sticky='news')

#------------------------------------First Person Buttons-------------------------------------------------------------------------------------
first_person_button_frame = tk.Frame(first_person_frame)

first_person_generate_button = tk.Button(first_person_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startFirstPersonThread)
first_person_back_button = tk.Button(first_person_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack buttons into button frame
first_person_generate_button.pack(side="left", padx=35)
first_person_back_button.pack(side="left", padx=30)

# place button frame on <something>
first_person_button_frame.pack(side="bottom", pady=10)

#--------------------------------us_uk_spellings Frame-----------------------------------------------------------------------------------------------
#--------------------------------------------------------------------------------------------------------------------------------------------------
us_uk_spellings_frame = tk.Frame(root, width=1000, height=600)
us_uk_spellings_frame.pack_propagate(False)
us_uk_spellings_frame.grid(row=0, column=0, sticky='news')

#--------------------------------us_uk_spellings Buttons-----------------------------------------------------------------------------------------------
# button frame - 
us_uk_spellings_button_frame = tk.Frame(us_uk_spellings_frame)
# US/UK Spellings button - based on dress IDs
us_uk_spellings_button = tk.Button(us_uk_spellings_button_frame, text="US/UK Spellings", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startUS_UK_SpellingsThread)

# back button
us_uk_spellings_back_button = tk.Button(us_uk_spellings_button_frame, text="Back", font=LABEL_FONT, width=15, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
us_uk_spellings_back_button.pack(side="right", padx=30)
# pack buttons into button frame - us_uk_spellings
us_uk_spellings_button.pack(side="left", padx=35)

# place button frame on us uk spellings 
us_uk_spellings_button_frame.pack(side="bottom", pady=10)

#--------------------------------ids_of_us_uk_spellings Frame-----------------------------------------------------------------------------------------------
#--------------------------------------------------------------------------------------------------------------------------------------------------
ids_of_us_uk_spellings_frame = tk.Frame(root, width=1000, height=600)
ids_of_us_uk_spellings_frame.pack_propagate(False)
ids_of_us_uk_spellings_frame.grid(row=0, column=0, sticky='news')

#--------------------------------ids_of_us_uk_spellings Buttons-----------------------------------------------------------------------------------------------
# button frame - 
ids_of_us_uk_spellings_button_frame = tk.Frame(ids_of_us_uk_spellings_frame)
# IDS OF US/UK Spellings button
ids_of_us_uk_spellings_button = tk.Button(ids_of_us_uk_spellings_button_frame, text="IDS OF US/UK Spellings", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startIDs_Of_US_UK_SpellingsThread)

# back button
ids_of_us_uk_spellings_back_button = tk.Button(ids_of_us_uk_spellings_button_frame, text="Back", font=LABEL_FONT, width=15, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
ids_of_us_uk_spellings_back_button.pack(side="right", padx=30)
# pack buttons into button frame - ids_of_us_uk_spellings
ids_of_us_uk_spellings_button.pack(side="left", padx=35)

# place button frame on us uk spellings 
ids_of_us_uk_spellings_button_frame.pack(side="bottom", pady=10)

'''
    ##Get Audio Frame
    =============================
'''

#--------------Get Audio Frame--------------------------------------------------------------------#
Get_audio_frame = tk.Frame(root, width=1000, height=600)
Get_audio_frame.pack_propagate(False)
Get_audio_frame.grid(row=0, column=0, sticky='news')

#button frame
Get_audio_button_frame = tk.Frame(Get_audio_frame)

#--------------------------------Insert the dress ID-----------------------------------------------------------------------------------------------
# Input text label 
text_field_label_ID_Address = tk.Label(root, text="Insert the dress ID", font=LABEL_FONT)
text_field_label_ID_Address.place(x=25, y=72.5)

# input text field
text_field_ID = tk.Entry(root)
text_field_ID.place(x=230, y=60, relwidth=.1, height=50)

text_field_ID = tk.Entry(root)
text_field_ID.place(x=230, y=60, relwidth=.1, height=50)

# #Get text
get_text_button = tk.Button(Get_audio_frame, text="get text", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=getTextThread)
get_text_button.place(x=350, y=60, relwidth=.1, height=50)

#Descript text label
text_field_label_Description = tk.Label(root, text="Description", font=LABEL_FONT)
text_field_label_Description.place(x=25, y=200)

# input text field
text_field_Description = tk.Text(Get_audio_frame)
text_field_Description.place(x=230, y=170, relwidth=.5, height=300)


#play Audio
play_audio_button = tk.Button(Get_audio_button_frame, text="Play Audio", font=LABEL_FONT, width=18, height=1, bg="#007FFF", fg="#ffffff", command=playAudio)
play_audio_button.pack(side="left", padx=30)

#Save Audio
save_audio_button = tk.Button(Get_audio_button_frame, text="Save Audio", font=LABEL_FONT, width=18, height=1, bg="#007FFF", fg="#ffffff", command=saveAudio)
save_audio_button.pack(side="left", padx=30)

# #upload Audio
upload_audio_button = tk.Button(Get_audio_button_frame, text="Upload Audio", font=LABEL_FONT, width=18, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
upload_audio_button.pack(side="left", padx=30)



# back button
audio_back_button = tk.Button(Get_audio_button_frame, text="Back", font=LABEL_FONT, width=18, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
audio_back_button.pack(side="left", padx=30)

# place button frame on diff report frame
Get_audio_button_frame.pack(side="bottom", pady=10)


'''
    ##Get All Audio Frame
    =============================
'''

#--------------------------------------------------------------------------------------------------------------------------------------------------
all_audio_frame = tk.Frame(root, width=1000, height=600)
all_audio_frame.pack_propagate(False)
all_audio_frame.grid(row=0, column=0, sticky='news')

#--------------------------------get all audios Buttons-----------------------------------------------------------------------------------------------
# button frame
all_audio_button_frame = tk.Frame(all_audio_frame)

# word analysis button
all_audio_button = tk.Button(all_audio_button_frame, text="Save audio", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=SaveAllAudios)
# back button
all_audio_back_button = tk.Button(all_audio_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
# pack buttons into button frame
all_audio_button.pack(side="left", padx=35)
all_audio_back_button.pack(side="left", padx=30)

# place button frame on word analysis frame
all_audio_button_frame.pack(side="bottom", pady=10)


'''
    ##DOB Analayser Frame
    =============================
'''

#--------------------------------------------------------------------------------------------------------------------------------------------------
DOB_Analyzer_frame = tk.Frame(root, width=1000, height=600)
DOB_Analyzer_frame.pack_propagate(False)
DOB_Analyzer_frame.grid(row=0, column=0, sticky='news')

#--------------------------------DOB Analayser Buttons-----------------------------------------------------------------------------------------------
# button frame
DOB_Analayser_button_frame = tk.Frame(DOB_Analyzer_frame)
data_frame = pd.DataFrame()  # Initialize an empty DataFrame

# # Create a Button to trigger the file upload
UPLOAD_FILE_button = tk.Button(root, text="Upload File", command=upload_file, bg="#007FFF", fg="#ffffff")
UPLOAD_FILE_button.place(x=500, y=60, relwidth=.1, height=50)

# DOB Analayser report button
DOB_Analayser_report_button = tk.Button(DOB_Analayser_button_frame, text="Generate HTML", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=generate_html)
# back button
DOB_Analayser_back_button = tk.Button(DOB_Analayser_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack buttons into button frame
DOB_Analayser_report_button.pack(side="left", padx=35)
DOB_Analayser_back_button.pack(side="left", padx=30)

# place button frame on DOB Analayser report frame
DOB_Analayser_button_frame.pack(side="bottom", pady=10)

#-------------------------------Word Search Frame-------------------------------------------------------------------------------------------
word_puzzle_frame = tk.Frame(root, width=1000, height=600)
word_puzzle_frame.pack_propagate(False)
word_puzzle_frame.grid(row=0, column=0, sticky='news')

#------------------------------------Word Search Buttons-------------------------------------------------------------------------------------
word_puzzle_button_frame = tk.Frame(word_puzzle_frame)

word_puzzle_generate_button = tk.Button(word_puzzle_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startWordPuzzleThread)
word_puzzle_back_button = tk.Button(word_puzzle_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack buttons into button frame
word_puzzle_generate_button.pack(side="left", padx=35)
word_puzzle_back_button.pack(side="left", padx=30)

# place button frame on <something>
word_puzzle_button_frame.pack(side="bottom", pady=10)

#--------------------------------Layout Radio Buttons-------------------------------------------------------------------------------------
# layout variable
layout = tk.IntVar()
layout.set(4)

# layout frame
layout_frame = tk.Frame(word_puzzle_frame)
# layout radio buttons
layout_radio4 = tk.Radiobutton(layout_frame, text="Puzzle on Left - Text on right - Single Page - Portrait Mode",
                                font=MAIN_FONT, variable=layout, value=4)
layout_radio1 = tk.Radiobutton(layout_frame, text="Puzzle on Left Page - Text on Right Page - Two Page Mode - Portrait Mode",
                                font=MAIN_FONT, variable=layout, value=1)
layout_radio2 = tk.Radiobutton(layout_frame, text="Puzzle on Right - Text on Left - Single Page - Landscape Mode",
                                font=MAIN_FONT, variable=layout, value=2)
layout_radio3 = tk.Radiobutton(layout_frame, text="Puzzle on Left - Text on Right - Single Page - Landscape Mode",
                                font=MAIN_FONT, variable=layout, value=3)
# pack radio buttons into layout frame
layout_radio4.pack(anchor="nw")
layout_radio1.pack(anchor="nw")
layout_radio2.pack(anchor="nw")
layout_radio3.pack(anchor="nw")
# place layout frame on main frame
layout_frame.place(x=175, y=150, width=800)
# layout radio buttons label
layout_label = tk.Label(word_puzzle_frame, text="Layout:", font=LABEL_FONT)
layout_label.place(x=25, y=170)

# separator line
separator1 = ttk.Separator(word_puzzle_frame)
separator1.place(x=175, y=150, relwidth=.8)

#--------------------------------Sort Radio Buttons---------------------------------------------------------------------------------------
# sort variable
sort_order = tk.IntVar()
sort_order.set(1)

# sort frame
sort_frame = tk.Frame(word_puzzle_frame)
# sort radio buttons
sort_radio1 = tk.Radiobutton(sort_frame, text="By Name", font=MAIN_FONT, variable=sort_order, value=1)
sort_radio2 = tk.Radiobutton(sort_frame, text="By ID", font=MAIN_FONT, variable=sort_order, value=2)
sort_radio3 = tk.Radiobutton(sort_frame, text="By Input Order", font=MAIN_FONT, variable=sort_order, value=3)
# pack radio buttons into sort frame
sort_radio1.pack(side="left")
sort_radio2.pack(side="left")
sort_radio3.pack(side="left")
# place sort frame on main frame
sort_frame.place(x=175, y=265, width=800)

# sort radio buttons label
sort_label = tk.Label(word_puzzle_frame, text="Sort Order:", font=LABEL_FONT)
sort_label.place(x=25, y=265)


# separator line
separator2 = ttk.Separator(word_puzzle_frame)
separator2.place(x=175, y=265, relwidth=.8)

#--------------------------------Word Radio Buttons---------------------------------------------------------------------------------------

#--------------------------------Preferences----------------------------------------------------------------------------------------------
preferences = {
    "WORD_COUNT": "10",
    "PUZ_WIDTH": "20",
}

# Using StringVar for dynamic updates
word_count_var = tk.StringVar(value=preferences["WORD_COUNT"])
puz_width_var = tk.StringVar(value=preferences["PUZ_WIDTH"])

# Frame for preference settings
preferences_frame = tk.Frame(word_puzzle_frame)
preferences_frame.place(x=175, y=295, width=800)

# Labels and entries for preferences
word_count_label = tk.Label(preferences_frame, text="Word Count:", font=MAIN_FONT)
word_count_label.grid(row=1, column=1, pady=10)
word_count = tk.Entry(preferences_frame, width=6, textvariable=word_count_var, font=MAIN_FONT)
word_count.grid(row=1, column=2)

puz_width_label = tk.Label(preferences_frame, text="Puzzle Width:", font=MAIN_FONT)
puz_width_label.grid(row=2, column=1, pady=10)
puz_width = tk.Entry(preferences_frame, width=6, textvariable=puz_width_var, font=MAIN_FONT)
puz_width.grid(row=2, column=2)


# Function to apply changes to preferences
def apply_changes():
    preferences["WORD_COUNT"] = word_count_var.get()
    preferences["PUZ_WIDTH"] = puz_width_var.get()
    print("Updated preferences:", preferences)

# Apply button
apply_button = tk.Button(preferences_frame, text="Apply Changes", command=apply_changes, font=LABEL_FONT)
apply_button.grid(row=4, column=1, columnspan=2, pady=10)

# preferences label
preferences_label = tk.Label(word_puzzle_frame, text="Preferences:", font=LABEL_FONT)
preferences_label.place(x=25, y=350)

# separator line
separator3 = ttk.Separator(word_puzzle_frame)
separator3.place(x=175, y=295, relwidth=.8)

#-------------------------------Start Main Frame----------------------------------------------------------------------------------------------

# raise main_frame to start
main_frame.tkraise()

# main gui loop
root.mainloop()