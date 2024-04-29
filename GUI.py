import sys, io, requests, threading, webbrowser, urllib, os, platform, openpyxl, string, textwrap, re, time, textstat
import traceback
import tkinter as tk
from tkinter import ttk, messagebox
from pptx import Presentation
import pptx.util
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import pandas as pd
from pandas import Series, DataFrame
from concurrent.futures import ThreadPoolExecutor, as_completed
from textblob import TextBlob
from bs4 import BeautifulSoup
import wikipediaapi
import openai
import random
from xml.dom.expatbuilder import FragmentBuilderNS
from localspelling import convert_spelling 
from localspelling.spelling_converter import get_dictionary 
import requests
import pyttsx3
import webbrowser
import tempfile
import datetime
from gtts import gTTS

from tkinter import filedialog
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.shapes.connector import Connector
from pptx.enum.shapes import MSO_SHAPE

from generate_pairs import startGeneratePairsThread
from us_uk_spelling import startIDs_Of_US_UK_SpellingsThread, startUS_UK_SpellingsThread
from diff_report import startDiffReportThread
from translate_package import startTranslationPackageThread
from get_and_save_audios import saveAudioThread, getTextThread, playAudioThread
from word_analysis import startWordAnalysisThread
from generate_wiki_link import startGenerateWikiLinkThread
# pip install googletrans
# pip install googletrans==4.0.0-rc1
import googletrans

openai.api_key = 'private'

ROOT_WIDTH = 1000 # app window width
ROOT_HEIGHT = 600 # app window height

root = tk.Tk()
root.title("Project ABCD Admin Panel")
sw_placement = int(root.winfo_screenwidth()/2 - ROOT_WIDTH/2) # to place at half width of screen
sh_placement = int(root.winfo_screenheight()/2 - ROOT_HEIGHT/2) # to place at half height of screen
root.geometry(f"{ROOT_WIDTH}x{ROOT_HEIGHT}+{sw_placement}+{sh_placement}")
root.minsize(ROOT_WIDTH, ROOT_HEIGHT)

# Set sys.stdout to use utf-8 encoding
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

MAIN_FONT = ("helvetica", 12)
LABEL_FONT = ("helvetica bold", 14)
LANGUAGES = [
    "Telugu",
    "Hindi",
    "Spanish"
]
# creates preferences dictionary from preferences.txt
# creates default preferences.txt file if one does not exist
try:
    with open("preferences.txt", "r", encoding="utf8") as file:
        lines = file.readlines()
        preferences = {}

        for line in lines:
            key, value = line.split('=')
            preferences[key.strip()] = value.strip().replace('“', '').replace('”', '').replace('"', '').replace("'", '')
except FileNotFoundError:
    preferences = {
        "TEXT_SIZE" : "14",
        "TEXT_FONT" : "Times New Roman",
        "TITLE_SIZE" : "32",
        "TITLE_FONT" : "Arial",
        "SUBTITLE_SIZE" : "24",
        "SUBTITLE_FONT" : "Arial",
        "PIC_WIDTH" : "720",
        "PIC_HEIGHT" : "1040"
    }
    tk.messagebox.showwarning(title='Warning', message='No preferences.txt file exists in directory. Default preferences.txt will be created and used.')
    print('No preferences.txt file exists in directory. Default preferences.txt will be created and used.')
    with open('preferences.txt', 'w') as f:
        f.write('TEXT_SIZE = 14\nTEXT_FONT = Times New Roman\nTITLE_SIZE = 32\nTITLE_FONT = Arial\nSUBTITLE_SIZE = 24\nSUBTITLE_FONT = Arial\nPIC_WIDTH = 720\nPIC_HEIGHT = 1040')

# pip install googletrans
# pip install googletrans==4.0.0-rc1
import googletrans

#--------------------------------Main Frame-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------------------------------------------------
# configure grid to fill extra space and center
tk.Grid.rowconfigure(root, 0, weight=1)
tk.Grid.columnconfigure(root, 0, weight=1)

# main frame
main_frame = tk.Frame(root, width=1000, height=1000)
main_frame.pack_propagate(False)
main_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Main Title-----------------------------------------------------------------------------------------------
# Create a label widget
title_label = tk.Label(main_frame, text="Project ABCD\nMain Menu",  font=('Arial', 20))
title_label.pack(pady=100)

#--------------------------------Main Buttons-----------------------------------------------------------------------------------------------
# Create buttons widget
## Button settings
main_button_frame = tk.Frame(main_frame)
main_button_frame.place(relx=.5, rely=.35, anchor='center')
button_width = 20
button_height = 3
button_bgd_color = "#007FFF"
button_font_color = "#ffffff"

## Generate Book: Gets selected dress from API and import into ppt
generate_book_button = tk.Button(main_button_frame, text="Generate Book", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('book_gen_frame'))
generate_book_button.pack(side="left", padx=50)

## Diff Report: Create a SQL file of dresses that got changed from excel sheet byt comparing to API
diff_report_button = tk.Button(main_button_frame, text="Difference Report", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('diff_report_frame'))
diff_report_button.pack(side="left", padx=50, anchor='center')

## Generate Book: Get selected dress that user input & put into a table (ID, Name, Description Count, DYK Count, Total Nouns Count, Total Adjectives Count)
word_analysis_report_button = tk.Button(main_button_frame, text="Word Analysis Report", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('word_analysis_frame'))
word_analysis_report_button.pack(side="left", padx=50)

main_button_frame2 = tk.Frame(main_frame)
main_button_frame2.place(relx=.5, rely=.50, anchor='center')

## Google Images: Create an Excel file with 3 image links to the selected dresses
google_image_button = tk.Button(main_button_frame2, text="Google Image", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('google_image_frame'))
google_image_button.pack(side="left", padx=50)

## Wiki Link: [FILL IN THE ACTION HERE]
wiki_link_button = tk.Button(main_button_frame2, text="Wiki Link", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('wiki_link_frame'))
wiki_link_button.pack(side="left", padx=50)

## My Pairs: Shows pairs when they are searched
who_are_my_pairs_button = tk.Button(main_button_frame2, text="Who Are My Pairs?", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('who_are_my_pairs_frame'))
who_are_my_pairs_button.pack(side="left", padx=50)

main_button_frame3 = tk.Frame(main_frame)
main_button_frame3.place(relx=.5, rely=.65, anchor='center')

## US Spelling: us_uk spellings
us_uk_spellings_button = tk.Button(main_button_frame3, text="US/UK Spellings", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('us_uk_spellings_frame'))
us_uk_spellings_button.pack(side="left", padx=50)

## IDS OF US Spelling: ids of us_uk spellings
ids_of_us_uk_spellings_button = tk.Button(main_button_frame3, text="IDS OF US/UK Spellings", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('ids_of_us_uk_spellings_frame'))
ids_of_us_uk_spellings_button.pack(side="left", padx=50)

## Generate Book: fetches English text from Api, uses google translate to generate "telugu" text, then creates HTML package with english text on page, and "telugu" text on another.
translation_package_button = tk.Button(main_button_frame3, text="Translation Package", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('translation_package_frame'))
translation_package_button.pack(side="left", padx=50)

main_button_frame4 = tk.Frame(main_frame)
main_button_frame4.place(relx=.5, rely=.80, anchor='center')

## First Person: fetches text from api, uses ChatGPT to reword the description and did you know text to first person
first_person_button = tk.Button(main_button_frame4, text="First Person Conversion", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('first_person_frame'))
first_person_button.pack(side="left", padx=50)

## Word Puzzle: generates and creates crossword puzzles based of words in character descriptions
word_puzzle_button = tk.Button(main_button_frame4, text="Word Puzzle Creator", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('word_puzzle_frame'))
word_puzzle_button.pack(side="left", padx=50)

main_button_frame5 = tk.Frame(main_frame)
main_button_frame5.place(relx =.5, rely=.95, anchor='center')

##Get Audio Frame
get_audio_button = tk.Button(main_button_frame5, text="Get Audio", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('Get_audio_frame'))
get_audio_button.pack(side="left", padx=50)

## get all Audio Frame
get_all_audio_button = tk.Button(main_button_frame5, text="Get All Audio", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('all_audio_frame'))
get_all_audio_button.pack(side="left", padx=50)

## get DOB Analyzer Frame
DOB_Analyzer_button = tk.Button(main_button_frame4, text="DOB Analyzer", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('DOB_Analyzer_frame'))
DOB_Analyzer_button.pack(side="left", padx=50)

#--------------------------------Book Gen Frame---------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------------------------------------------------------
# book gen frame
book_gen_frame = tk.Frame(root, width=1000, height=600)
book_gen_frame.pack_propagate(False)
book_gen_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Text Field-----------------------------------------------------------------------------------------------
# text field label
text_field_label = tk.Label(root, text="Dress Numbers:", font=LABEL_FONT)
text_field_label.place(x=25, y=72.5)

# text field
text_field = tk.Text(root)
text_field.place(x=175, y=10, relwidth=.8, height=135)

# text field initialization
try:
    with open('slide_numbers.txt', 'r') as file:
        slide_number_content = file.readline().strip()
        text_field.insert("1.0", slide_number_content)
except FileNotFoundError:
    print(FileNotFoundError)

#--------------------------------Layout Radio Buttons-------------------------------------------------------------------------------------
# layout variable
layout = tk.IntVar()
layout.set(4)

# layout frame
layout_frame = tk.Frame(book_gen_frame)
# layout radio buttons
layout_radio4 = tk.Radiobutton(layout_frame, text="Picture on Left - Text on right - Single Page - Portrait Mode",
                                font=MAIN_FONT, variable=layout, value=4)
layout_radio1 = tk.Radiobutton(layout_frame, text="Picture on Left Page - Text on Right Page - Two Page Mode - Portrait Mode",
                                font=MAIN_FONT, variable=layout, value=1)
layout_radio2 = tk.Radiobutton(layout_frame, text="Picture on Right - Text on Left - Single Page - Landscape Mode",
                                font=MAIN_FONT, variable=layout, value=2)
layout_radio3 = tk.Radiobutton(layout_frame, text="Picture on Left - Text on Right - Single Page - Landscape Mode",
                                font=MAIN_FONT, variable=layout, value=3)
# pack radio buttons into layout frame
layout_radio4.pack(anchor="nw")
layout_radio1.pack(anchor="nw")
layout_radio2.pack(anchor="nw")
layout_radio3.pack(anchor="nw")
# place layout frame on main frame
layout_frame.place(x=175, y=150, width=800)
# layout radio buttons label
layout_label = tk.Label(book_gen_frame, text="Layout:", font=LABEL_FONT)
layout_label.place(x=25, y=170)

# separator line
separator1 = ttk.Separator(book_gen_frame)
separator1.place(x=175, y=150, relwidth=.8)

#--------------------------------Sort Radio Buttons---------------------------------------------------------------------------------------
# sort variable
sort_order = tk.IntVar()
sort_order.set(1)

# sort frame
sort_frame = tk.Frame(book_gen_frame)
# sort radio buttons
sort_radio1 = tk.Radiobutton(sort_frame, text="By Name", font=MAIN_FONT, variable=sort_order, value=1)
sort_radio2 = tk.Radiobutton(sort_frame, text="By ID", font=MAIN_FONT, variable=sort_order, value=2)
sort_radio3 = tk.Radiobutton(sort_frame, text="By Input Order", font=MAIN_FONT, variable=sort_order, value=3)
# pack radio buttons into sort frame
sort_radio1.pack(side="left")
sort_radio2.pack(side="left")
sort_radio3.pack(side="left")
# place sort frame on main frame
sort_frame.place(x=175, y=265, width=800)

# sort radio buttons label
sort_label = tk.Label(book_gen_frame, text="Sort Order:", font=LABEL_FONT)
sort_label.place(x=25, y=265)

# separator line
separator2 = ttk.Separator(book_gen_frame)
separator2.place(x=175, y=265, relwidth=.8)


#--------------------------------Preferences----------------------------------------------------------------------------------------------
# Initialize variables for storing values
text_size_var = tk.StringVar()
title_size_var = tk.StringVar()
subtitle_size_var = tk.StringVar()
text_font_var = tk.StringVar()
title_font_var = tk.StringVar()
subtitle_font_var = tk.StringVar()
pic_width_var = tk.StringVar()
pic_height_var = tk.StringVar()

# Set initial values for the Entry fields
text_size_var.set(preferences["TEXT_SIZE"])
title_size_var.set(preferences["TITLE_SIZE"])
subtitle_size_var.set(preferences["SUBTITLE_SIZE"])
text_font_var.set(preferences["TEXT_FONT"])
title_font_var.set(preferences["TITLE_FONT"])
subtitle_font_var.set(preferences["SUBTITLE_FONT"])
pic_width_var.set(preferences["PIC_WIDTH"])
pic_height_var.set(preferences["PIC_HEIGHT"])

# preferences frame
preferences_frame = tk.Frame(book_gen_frame)
# preferences labels and entry fields
text_size_label = tk.Label(preferences_frame, text="Text Size:", font=MAIN_FONT)
text_size = tk.Entry(preferences_frame, width=3, textvariable=text_size_var, state="disabled", font=MAIN_FONT)

title_size_label = tk.Label(preferences_frame, text="Title Size:", font=MAIN_FONT)
title_size = tk.Entry(preferences_frame, width=3, textvariable=title_size_var, state="disabled",  font=MAIN_FONT)

subtitle_size_label = tk.Label(preferences_frame, text="Subtitle Size:", font=MAIN_FONT)
subtitle_size = tk.Entry(preferences_frame, width=3, textvariable=subtitle_size_var, state="disabled",  font=MAIN_FONT)

text_font_label = tk.Label(preferences_frame, text="Text Font:", font=MAIN_FONT)
text_font = tk.Entry(preferences_frame, width=25, textvariable=text_font_var, state="disabled",  font=MAIN_FONT)

title_font_label = tk.Label(preferences_frame, text="Title Font:", font=MAIN_FONT)
title_font = tk.Entry(preferences_frame, width=25, textvariable=title_font_var, state="disabled",  font=MAIN_FONT)

subtitle_font_label = tk.Label(preferences_frame, text="Subitle Font:", font=MAIN_FONT)
subtitle_font = tk.Entry(preferences_frame, width=25, textvariable=subtitle_font_var, state="disabled",  font=MAIN_FONT)

pic_width_label = tk.Label(preferences_frame, text="Pic Width:", font=MAIN_FONT)
pic_width = tk.Entry(preferences_frame, width=6, textvariable=pic_width_var, state="disabled",  font=MAIN_FONT)

pic_height_label = tk.Label(preferences_frame, text="Pic Height:", font=MAIN_FONT)
pic_height = tk.Entry(preferences_frame, width=6, textvariable=pic_height_var, state="disabled",  font=MAIN_FONT)

# grid preference labels and entry fields into preferences frame
# column 1 + 2
text_size_label.grid(row=1, column=1, pady=10)
text_size.grid(row=1, column=2)

title_size_label.grid(row=2, column=1, pady=10)
title_size.grid(row=2, column=2)

subtitle_size_label.grid(row=3, column=1, pady=10, padx=15)
subtitle_size.grid(row=3, column=2)

# column 3 + 4
text_font_label.grid(row=1, column=3)
text_font.grid(row=1, column=4)

title_font_label.grid(row=2, column=3)
title_font.grid(row=2, column=4)

subtitle_font_label.grid(row=3, column=3, padx=15)
subtitle_font.grid(row=3, column=4)

# column 5 + 6
pic_width_label.grid(row=1, column=5)
pic_width.grid(row=1, column=6)

pic_height_label.grid(row=2, column=5, padx=15)
pic_height.grid(row=2, column=6)

# place preferences frame on main frame
preferences_frame.place(x=175, y=295, width=800)

# preferences label
preferences_label = tk.Label(book_gen_frame, text="Preferences:", font=LABEL_FONT)
preferences_label.place(x=25, y=350)

# separator line
separator3 = ttk.Separator(book_gen_frame)
separator3.place(x=175, y=295, relwidth=.8)

#--------------------------------Numbering Radio Buttons--------------------------------------------------------------------------------------
# numbering variable
numbering = tk.IntVar()
numbering.set(1)

# numbering frame
numbering_frame = tk.Frame(book_gen_frame)
# numbering radio buttons
numbering_radio1 = tk.Radiobutton(numbering_frame, text="Show both Page Number and Dress ID", font=MAIN_FONT, variable=numbering, value=1)
numbering_radio2 = tk.Radiobutton(numbering_frame, text="Show Page Number", font=MAIN_FONT, variable=numbering, value=2)
numbering_radio3 = tk.Radiobutton(numbering_frame, text="Show Dress ID", font=MAIN_FONT, variable=numbering, value=3)
# pack numbering buttons into sort frame
numbering_radio1.pack(side="left")
numbering_radio2.pack(side="left")
numbering_radio3.pack(side="left")
# place numbering frame on main frame
numbering_frame.place(x=175, y=445, width=800)

# numbering radio buttons label
numbering_label = tk.Label(book_gen_frame, text="Numbering:", font=LABEL_FONT)
numbering_label.place(x=25, y=445)

# separator line
separator4 = ttk.Separator(book_gen_frame)
separator4.place(x=175, y=445, relwidth=.8)

#--------------------------------Translate and Image Check Buttons--------------------------------------------------------------------------------------
# translate check variable
translate = tk.IntVar()
translate.set(0)
# language options variable
language = tk.StringVar()
language.set(LANGUAGES[0])
# Download image variable
download_imgs = tk.IntVar()
download_imgs.set(0)
# Generate book from local Excel sheet
gen_local = tk.IntVar()
gen_local.set(0)

# translate frame
check_button_frame = tk.Frame(book_gen_frame)
# translate check button
translate_checkbutton = tk.Checkbutton(check_button_frame, text="Translate to:", font=MAIN_FONT, variable=translate, onvalue=1, offvalue=0)
# language options
language_options = tk.OptionMenu(check_button_frame, language, *LANGUAGES)
# download images
download_images = tk.Checkbutton(check_button_frame, text="Download Images", font=MAIN_FONT, variable=download_imgs, onvalue=1, offvalue=0, command=lambda: gen_local.set(0))
# generate book from local Excel sheet
generate_from_local = tk.Checkbutton(check_button_frame, text="Generate from Local", font=MAIN_FONT, variable=gen_local, onvalue=1, offvalue=0, command=lambda: download_imgs.set(0))

# pack translate options into translate frame
translate_checkbutton.pack(side="left")
language_options.pack(side="left")
download_images.pack(side="left")
generate_from_local.pack(side="left")
# place translate frame on main frame
check_button_frame.place(x=175, y=495)

# separator line
separator5 = ttk.Separator(book_gen_frame)
separator5.place(x=175, y=495, relwidth=.8)

#--------------------------------Book Gen Buttons--------------------------------------------------------------------------------------
# button frame
book_gen_button_frame = tk.Frame(book_gen_frame)
# generate button
book_gen_generate_button = tk.Button(book_gen_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff")
# help button
book_gen_help_button = tk.Button(book_gen_button_frame, text="Help", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff")
# upload button
book_gen_back_button = tk.Button(book_gen_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
# pack buttons into button frame
book_gen_generate_button.pack(side="left", padx=35)
book_gen_help_button.pack(side="left")
book_gen_back_button.pack(side="left", padx=30)

# place button frame on main frame
book_gen_button_frame.pack(side="bottom", pady=10)


#--------------------------------Diff Report Frame-----------------------------------------------------------------------------------------------
#------------------------------------------------------------------------------------------------------------------------------------------------
diff_report_frame = tk.Frame(root, width=1000, height=600)
diff_report_frame.pack_propagate(False)
diff_report_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Diff Report Buttons-----------------------------------------------------------------------------------------------
# button frame
diff_report_button_frame = tk.Frame(diff_report_frame)
# difference report button
diff_report_button = tk.Button(diff_report_button_frame, text="Diff Report", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: startDiffReportThread(root, text_field, diff_report_button))
# back button
diff_back_button = tk.Button(diff_report_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
# pack buttons into button frame
diff_report_button.pack(side="left", padx=35)
diff_back_button.pack(side="left", padx=30)

# place button frame on diff report frame
diff_report_button_frame.pack(side="bottom", pady=10)


#--------------------------------Word Analysis Frame-----------------------------------------------------------------------------------------------
#--------------------------------------------------------------------------------------------------------------------------------------------------
word_analysis_frame = tk.Frame(root, width=1000, height=600)
word_analysis_frame.pack_propagate(False)
word_analysis_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Word Analysis Buttons-----------------------------------------------------------------------------------------------
# button frame
word_analysis_button_frame = tk.Frame(word_analysis_frame)
# word analysis button
word_analysis_button = tk.Button(word_analysis_button_frame, text="Word Analysis", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: startWordAnalysisThread(root, text_field, word_analysis_button))
# back button
word_analysis_back_button = tk.Button(word_analysis_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
# pack buttons into button frame
word_analysis_button.pack(side="left", padx=35)
word_analysis_back_button.pack(side="left", padx=30)

# place button frame on word analysis frame
word_analysis_button_frame.pack(side="bottom", pady=10)

#--------------------------------Google Image Frame-----------------------------------------------------------------------------------------------
#--------------------------------------------------------------------------------------------------------------------------------------------------
google_image_frame = tk.Frame(root, width=1000, height=600)
google_image_frame.pack_propagate(False)
google_image_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Google Image Buttons-----------------------------------------------------------------------------------------------
# Download google image variable
download_google_imgs = tk.IntVar()
download_google_imgs.set(0)

# button frame
google_image_button_frame = tk.Frame(google_image_frame)
# checkbox frame
download_google_image_checkbutton = tk.Frame(google_image_frame)
# download images checkbox
download_google_images = tk.Checkbutton(download_google_image_checkbutton, text="Download Images", font=MAIN_FONT, variable=download_google_imgs, onvalue=1, offvalue=0, command=lambda: gen_local.set(0))
# google search button
google_image_search_button = tk.Button(google_image_button_frame, text="Google Search", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff")
# back button
google_image_back_button = tk.Button(google_image_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack buttons into button frame
download_google_image_checkbutton.pack(side="left")
download_google_images.pack(side="left")
google_image_search_button.pack(side="left", padx=35)
google_image_back_button.pack(side="left", padx=30)
# place google image frame on main frame
google_image_button_frame.place(x=175, y=495)
download_google_image_checkbutton.place(x=170, y=150)

# place button frame on word analysis frame
google_image_button_frame.pack(side="bottom", pady=10)

#--------------------------------Wiki Link Frame-----------------------------------------------------------------------------------------------
#------------------------------------------------------------------------------------------------------------------------------------------------
wiki_link_frame = tk.Frame(root, width=1000, height=600)
wiki_link_frame.pack_propagate(False)
wiki_link_frame.grid(row=0, column=0, sticky='news')
 
#--------------------------------Wiki Link Buttons--------------------------------------------------------------------------------------
# button frame
wiki_link_gen_button_frame = tk.Frame(wiki_link_frame)

wiki_link_gen_button = tk.Button(wiki_link_gen_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: startGenerateWikiLinkThread(root, text_field, wiki_link_gen_button))
wiki_link_back_button = tk.Button(wiki_link_gen_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack buttons into button frame
wiki_link_gen_button.pack(side="left", padx=35)
wiki_link_back_button.pack(side="left", padx=30)

# place button frame on word analysis frame
wiki_link_gen_button_frame.pack(side="bottom", pady=10)

#--------------------------------Who Are My Pairs Frame-----------------------------------------------------------------------------------------------
#------------------------------------------------------------------------------------------------------------------------------------------------
who_are_my_pairs_frame = tk.Frame(root, width=1000, height=600)
who_are_my_pairs_frame.pack_propagate(False)
who_are_my_pairs_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Who Are My Pairs Buttons-----------------------------------------------------------------------------------
#button frame
who_are_my_pairs_gen_button_frame = tk.Frame(who_are_my_pairs_frame)

who_are_my_pairs_gen_button = tk.Button(who_are_my_pairs_gen_button_frame, text="Generate Pairs", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: startGeneratePairsThread(root, text_field, who_are_my_pairs_gen_button))
who_are_my_pairs_back_button = tk.Button(who_are_my_pairs_gen_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack button into button frame
who_are_my_pairs_gen_button.pack(side="left", padx=35)
who_are_my_pairs_back_button.pack(side="left", padx=30)

# place button frame on word analysis frame?
who_are_my_pairs_gen_button_frame.pack(side="bottom", pady=10)

#--------------------------------Translation Package Frame-----------------------------------------------------------------------------------------------
#--------------------------------------------------------------------------------------------------------------------------------------------------
translation_package_frame = tk.Frame(root, width=1000, height=600)
translation_package_frame.pack_propagate(False)
translation_package_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Translation Package Buttons-----------------------------------------------------------------------------------------------
# generate button frame
translation_package_button_frame = tk.Frame(translation_package_frame)
# generate button
translation_package_generate_button = tk.Button(translation_package_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: startTranslationPackageThread(root, text_field, translation_package_generate_button))
# help button
translation_package_help_button = tk.Button(translation_package_button_frame, text="Help", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff")
# upload button
translation_package_back_button = tk.Button(translation_package_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
# pack buttons into button frame
translation_package_generate_button.pack(side="left", padx=35)
translation_package_help_button.pack(side="left")
translation_package_back_button.pack(side="left", padx=30)

# place button frame on word analysis frame
translation_package_button_frame.pack(side="bottom", pady=10)

#-------------------------------First Person Frame-------------------------------------------------------------------------------------------
first_person_frame = tk.Frame(root, width=1000, height=600)
first_person_frame.pack_propagate(False)
first_person_frame.grid(row=0, column=0, sticky='news')

#------------------------------------First Person Buttons-------------------------------------------------------------------------------------
first_person_button_frame = tk.Frame(first_person_frame)

first_person_generate_button = tk.Button(first_person_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff")
first_person_back_button = tk.Button(first_person_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack buttons into button frame
first_person_generate_button.pack(side="left", padx=35)
first_person_back_button.pack(side="left", padx=30)

# place button frame on <something>
first_person_button_frame.pack(side="bottom", pady=10)

#--------------------------------us_uk_spellings Frame-----------------------------------------------------------------------------------------------
#--------------------------------------------------------------------------------------------------------------------------------------------------
us_uk_spellings_frame = tk.Frame(root, width=1000, height=600)
us_uk_spellings_frame.pack_propagate(False)
us_uk_spellings_frame.grid(row=0, column=0, sticky='news')

#--------------------------------us_uk_spellings Buttons-----------------------------------------------------------------------------------------------
# button frame - 
us_uk_spellings_button_frame = tk.Frame(us_uk_spellings_frame)
# US/UK Spellings button - based on dress IDs
us_uk_spellings_button = tk.Button(us_uk_spellings_button_frame, text="US/UK Spellings", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: startUS_UK_SpellingsThread(root, text_field, us_uk_spellings_button))

# back button
us_uk_spellings_back_button = tk.Button(us_uk_spellings_button_frame, text="Back", font=LABEL_FONT, width=15, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
us_uk_spellings_back_button.pack(side="right", padx=30)
# pack buttons into button frame - us_uk_spellings
us_uk_spellings_button.pack(side="left", padx=35)

# place button frame on us uk spellings 
us_uk_spellings_button_frame.pack(side="bottom", pady=10)

#--------------------------------ids_of_us_uk_spellings Frame-----------------------------------------------------------------------------------------------
#--------------------------------------------------------------------------------------------------------------------------------------------------
ids_of_us_uk_spellings_frame = tk.Frame(root, width=1000, height=600)
ids_of_us_uk_spellings_frame.pack_propagate(False)
ids_of_us_uk_spellings_frame.grid(row=0, column=0, sticky='news')

#--------------------------------ids_of_us_uk_spellings Buttons-----------------------------------------------------------------------------------------------
# button frame - 
ids_of_us_uk_spellings_button_frame = tk.Frame(ids_of_us_uk_spellings_frame)
# IDS OF US/UK Spellings button
ids_of_us_uk_spellings_button = tk.Button(ids_of_us_uk_spellings_button_frame, text="IDS OF US/UK Spellings", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: startIDs_Of_US_UK_SpellingsThread(root, text_field, ids_of_us_uk_spellings_button))

# back button
ids_of_us_uk_spellings_back_button = tk.Button(ids_of_us_uk_spellings_button_frame, text="Back", font=LABEL_FONT, width=15, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
ids_of_us_uk_spellings_back_button.pack(side="right", padx=30)
# pack buttons into button frame - ids_of_us_uk_spellings
ids_of_us_uk_spellings_button.pack(side="left", padx=35)

# place button frame on us uk spellings 
ids_of_us_uk_spellings_button_frame.pack(side="bottom", pady=10)

'''
    ##Get Audio Frame
    =============================
'''

#--------------Get Audio Frame--------------------------------------------------------------------#
Get_audio_frame = tk.Frame(root, width=1000, height=600)
Get_audio_frame.pack_propagate(False)
Get_audio_frame.grid(row=0, column=0, sticky='news')

#button frame
Get_audio_button_frame = tk.Frame(Get_audio_frame)

#--------------------------------Insert the dress ID-----------------------------------------------------------------------------------------------
# Input text label 
text_field_label_ID_Address = tk.Label(root, text="Insert the dress ID", font=LABEL_FONT)
text_field_label_ID_Address.place(x=25, y=72.5)

# input text field
text_field_ID = tk.Entry(root)
text_field_ID.place(x=230, y=60, relwidth=.1, height=50)

text_field_ID = tk.Entry(root)
text_field_ID.place(x=230, y=60, relwidth=.1, height=50)

# #Get text
get_text_button = tk.Button(Get_audio_frame, text="get text", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: getTextThread(root, text_field_Description, text_field_ID))
get_text_button.place(x=350, y=60, relwidth=.1, height=50)

#Descript text label
text_field_label_Description = tk.Label(root, text="Description", font=LABEL_FONT)
text_field_label_Description.place(x=25, y=200)

# input text field
text_field_Description = tk.Text(Get_audio_frame)
text_field_Description.place(x=230, y=170, relwidth=.5, height=300)


#play Audio
play_audio_button = tk.Button(Get_audio_button_frame, text="Play Audio", font=LABEL_FONT, width=18, height=1, bg="#007FFF", fg="#ffffff", command=lambda: playAudioThread(root, play_audio_button, text_field_Description))
play_audio_button.pack(side="left", padx=30)

#Save Audio
save_audio_button = tk.Button(Get_audio_button_frame, text="Save Audio", font=LABEL_FONT, width=18, height=1, bg="#007FFF", fg="#ffffff", command=lambda: saveAudioThread(root, text_field_Description, save_audio_button, text_field_ID))
save_audio_button.pack(side="left", padx=30)

# #upload Audio
upload_audio_button = tk.Button(Get_audio_button_frame, text="Upload Audio", font=LABEL_FONT, width=18, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
upload_audio_button.pack(side="left", padx=30)



# back button
audio_back_button = tk.Button(Get_audio_button_frame, text="Back", font=LABEL_FONT, width=18, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
audio_back_button.pack(side="left", padx=30)

# place button frame on diff report frame
Get_audio_button_frame.pack(side="bottom", pady=10)


'''
    ##Get All Audio Frame
    =============================
'''

#--------------------------------------------------------------------------------------------------------------------------------------------------
all_audio_frame = tk.Frame(root, width=1000, height=600)
all_audio_frame.pack_propagate(False)
all_audio_frame.grid(row=0, column=0, sticky='news')

#--------------------------------get all audios Buttons-----------------------------------------------------------------------------------------------
# button frame
all_audio_button_frame = tk.Frame(all_audio_frame)

# word analysis button
all_audio_button = tk.Button(all_audio_button_frame, text="Save audio", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff")
# back button
all_audio_back_button = tk.Button(all_audio_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
# pack buttons into button frame
all_audio_button.pack(side="left", padx=35)
all_audio_back_button.pack(side="left", padx=30)

# place button frame on word analysis frame
all_audio_button_frame.pack(side="bottom", pady=10)


'''
    ##DOB Analayser Frame
    =============================
'''

#--------------------------------------------------------------------------------------------------------------------------------------------------
DOB_Analyzer_frame = tk.Frame(root, width=1000, height=600)
DOB_Analyzer_frame.pack_propagate(False)
DOB_Analyzer_frame.grid(row=0, column=0, sticky='news')

#--------------------------------DOB Analayser Buttons-----------------------------------------------------------------------------------------------
# button frame
DOB_Analayser_button_frame = tk.Frame(DOB_Analyzer_frame)
data_frame = pd.DataFrame()  # Initialize an empty DataFrame

# # Create a Button to trigger the file upload
UPLOAD_FILE_button = tk.Button(root, text="Upload File", bg="#007FFF", fg="#ffffff")
UPLOAD_FILE_button.place(x=500, y=60, relwidth=.1, height=50)

# DOB Analayser report button
DOB_Analayser_report_button = tk.Button(DOB_Analayser_button_frame, text="Generate HTML", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff")
# back button
DOB_Analayser_back_button = tk.Button(DOB_Analayser_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack buttons into button frame
DOB_Analayser_report_button.pack(side="left", padx=35)
DOB_Analayser_back_button.pack(side="left", padx=30)

# place button frame on DOB Analayser report frame
DOB_Analayser_button_frame.pack(side="bottom", pady=10)

#-------------------------------Word Search Frame-------------------------------------------------------------------------------------------
word_puzzle_frame = tk.Frame(root, width=1000, height=600)
word_puzzle_frame.pack_propagate(False)
word_puzzle_frame.grid(row=0, column=0, sticky='news')

#------------------------------------Word Search Buttons-------------------------------------------------------------------------------------
word_puzzle_button_frame = tk.Frame(word_puzzle_frame)

word_puzzle_generate_button = tk.Button(word_puzzle_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff")
word_puzzle_back_button = tk.Button(word_puzzle_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack buttons into button frame
word_puzzle_generate_button.pack(side="left", padx=35)
word_puzzle_back_button.pack(side="left", padx=30)

# place button frame on <something>
word_puzzle_button_frame.pack(side="bottom", pady=10)

#--------------------------------Layout Radio Buttons-------------------------------------------------------------------------------------
# layout variable
layout = tk.IntVar()
layout.set(4)

# layout frame
layout_frame = tk.Frame(word_puzzle_frame)
# layout radio buttons
layout_radio4 = tk.Radiobutton(layout_frame, text="Puzzle on Left - Text on right - Single Page - Portrait Mode",
                                font=MAIN_FONT, variable=layout, value=4)
layout_radio1 = tk.Radiobutton(layout_frame, text="Puzzle on Left Page - Text on Right Page - Two Page Mode - Portrait Mode",
                                font=MAIN_FONT, variable=layout, value=1)
layout_radio2 = tk.Radiobutton(layout_frame, text="Puzzle on Right - Text on Left - Single Page - Landscape Mode",
                                font=MAIN_FONT, variable=layout, value=2)
layout_radio3 = tk.Radiobutton(layout_frame, text="Puzzle on Left - Text on Right - Single Page - Landscape Mode",
                                font=MAIN_FONT, variable=layout, value=3)
# pack radio buttons into layout frame
layout_radio4.pack(anchor="nw")
layout_radio1.pack(anchor="nw")
layout_radio2.pack(anchor="nw")
layout_radio3.pack(anchor="nw")
# place layout frame on main frame
layout_frame.place(x=175, y=150, width=800)
# layout radio buttons label
layout_label = tk.Label(word_puzzle_frame, text="Layout:", font=LABEL_FONT)
layout_label.place(x=25, y=170)

# separator line
separator1 = ttk.Separator(word_puzzle_frame)
separator1.place(x=175, y=150, relwidth=.8)

#--------------------------------Sort Radio Buttons---------------------------------------------------------------------------------------
# sort variable
sort_order = tk.IntVar()
sort_order.set(1)

# sort frame
sort_frame = tk.Frame(word_puzzle_frame)
# sort radio buttons
sort_radio1 = tk.Radiobutton(sort_frame, text="By Name", font=MAIN_FONT, variable=sort_order, value=1)
sort_radio2 = tk.Radiobutton(sort_frame, text="By ID", font=MAIN_FONT, variable=sort_order, value=2)
sort_radio3 = tk.Radiobutton(sort_frame, text="By Input Order", font=MAIN_FONT, variable=sort_order, value=3)
# pack radio buttons into sort frame
sort_radio1.pack(side="left")
sort_radio2.pack(side="left")
sort_radio3.pack(side="left")
# place sort frame on main frame
sort_frame.place(x=175, y=265, width=800)

# sort radio buttons label
sort_label = tk.Label(word_puzzle_frame, text="Sort Order:", font=LABEL_FONT)
sort_label.place(x=25, y=265)


# separator line
separator2 = ttk.Separator(word_puzzle_frame)
separator2.place(x=175, y=265, relwidth=.8)

#--------------------------------Word Radio Buttons---------------------------------------------------------------------------------------

#--------------------------------Preferences----------------------------------------------------------------------------------------------
preferences = {
    "WORD_COUNT": "10",
    "PUZ_WIDTH": "20",
}

# Using StringVar for dynamic updates
word_count_var = tk.StringVar(value=preferences["WORD_COUNT"])
puz_width_var = tk.StringVar(value=preferences["PUZ_WIDTH"])

# Frame for preference settings
preferences_frame = tk.Frame(word_puzzle_frame)
preferences_frame.place(x=175, y=295, width=800)

# Labels and entries for preferences
word_count_label = tk.Label(preferences_frame, text="Word Count:", font=MAIN_FONT)
word_count_label.grid(row=1, column=1, pady=10)
word_count = tk.Entry(preferences_frame, width=6, textvariable=word_count_var, font=MAIN_FONT)
word_count.grid(row=1, column=2)

puz_width_label = tk.Label(preferences_frame, text="Puzzle Width:", font=MAIN_FONT)
puz_width_label.grid(row=2, column=1, pady=10)
puz_width = tk.Entry(preferences_frame, width=6, textvariable=puz_width_var, font=MAIN_FONT)
puz_width.grid(row=2, column=2)


# Function to apply changes to preferences
def apply_changes():
    preferences["WORD_COUNT"] = word_count_var.get()
    preferences["PUZ_WIDTH"] = puz_width_var.get()
    print("Updated preferences:", preferences)

# Apply button
apply_button = tk.Button(preferences_frame, text="Apply Changes", command=apply_changes, font=LABEL_FONT)
apply_button.grid(row=4, column=1, columnspan=2, pady=10)

# preferences label
preferences_label = tk.Label(word_puzzle_frame, text="Preferences:", font=LABEL_FONT)
preferences_label.place(x=25, y=350)

# separator line
separator3 = ttk.Separator(word_puzzle_frame)
separator3.place(x=175, y=295, relwidth=.8)

'''
Gathers data from API
'''
def downloadAPIData(url, id_number):
    try:
        response = requests.get(url, headers={"User-Agent": "XY"})
        # append dress info to dress data if response status_code == 200
        if response.ok:
            return response.json()['data']
        else:
            print(f'Request for dress ID: {id_number} failed.')
    except requests.exceptions.RequestException as e:
        print(f'-- DEBUG -- in downloadAPIData: {e}')
    except Exception as e:
        # tk.messagebox.showerror(title="Error", message=f'Could not make connection!\n\nError: {e}')
        print(f'Error: {e}')

def raiseFrame(frame):
    if frame == 'main_frame':
        main_frame.tkraise()
        root.title("Project ABCD Admin Panel")
    elif frame == 'book_gen_frame':
        book_gen_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Book Generation")
    elif frame == 'diff_report_frame':
        diff_report_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Difference Report")
    elif frame == 'word_analysis_frame':
        word_analysis_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Word Analysis")
    elif frame == 'google_image_frame':
        google_image_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Google Image")
    elif frame == 'wiki_link_frame':
        wiki_link_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Wiki Link")
    elif frame == 'who_are_my_pairs_frame':
        who_are_my_pairs_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Who Are My Pairs")
    elif frame == 'translation_package_frame':  
        translation_package_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Translation Package")
    elif frame == 'first_person_frame':
        first_person_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
    elif frame == 'us_uk_spellings_frame':
        us_uk_spellings_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD US/UK Spellings")
    elif frame == 'ids_of_us_uk_spellings_frame':
        ids_of_us_uk_spellings_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD IDs_Of_US/UK Spellings")
    elif frame == 'Get_audio_frame':
        Get_audio_frame.tkraise()
        text_field_label_ID_Address.tkraise()
        # text_field_label_ID_Address.tkraise()
        text_field_label_Description.tkraise()
        text_field_Description.tkraise()
        text_field_ID.tkraise()
        play_audio_button.tkraise()
        save_audio_button.tkraise()
        upload_audio_button.tkraise()
        get_text_button.tkraise()
        root.title("Project ABCD Get Audio")
    elif frame == "all_audio_frame":
        all_audio_frame.tkraise()
        all_audio_button.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Word Analysis")
    elif frame == "DOB_Analyzer_frame":
        DOB_Analyzer_frame.tkraise()
        word_analysis_button.tkraise()
        word_analysis_back_button.tkraise()
        UPLOAD_FILE_button.tkraise()
        # UPLOAD_FILE_button.tkraise()
        root.title("Project ABCD Word Analysis")
    elif frame == 'word_puzzle_frame':
        word_puzzle_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Word Search Puzzles")
#-------------------------------Start Main Frame----------------------------------------------------------------------------------------------

# raise main_frame to start
main_frame.tkraise()

# main gui loop
root.mainloop()